#!/usr/bin/env python3
"""AIエージェントの進化 ― RPAの次に来るもの - PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xC7)
COLOR_ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: タイトル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "AIエージェントの進化", font_size=48, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
            "RPAの次に来るもの ― 企業が本当に必要とするAI自動化とは",
            font_size=22, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.2), Inches(3.3), COLOR_ACCENT)

add_textbox(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
            "― 汎用AIツールから業務特化型AIオフィスへ ―",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 従来型自動化（RPA）の限界
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "従来型自動化（RPA）の限界", font_size=36, bold=True)

problems = [
    ("脆弱性", "UI変更で\nシナリオが壊れる", COLOR_RED),
    ("高コスト", "開発・保守に\n専門人材が必要", COLOR_ORANGE),
    ("柔軟性の欠如", "例外処理に\n対応できない", COLOR_RED),
    ("スケールの壁", "業務が増えるほど\nシナリオも増加", COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(3.5), RGBColor(0x24, 0x24, 0x3E), 0.05)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.4), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, x + Inches(0.9), y + Inches(0.55), Inches(1.0), Inches(0.7),
                "!" if color == COLOR_RED else "?", font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.4), Inches(0.6),
                title, font_size=24, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(1.0),
                desc, font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "RPAは「画面操作の記録・再生」であり、業務の本質的な理解には至らない",
            font_size=15, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: AIエージェントとは何か ― 4つの基本機能
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "AIエージェントの4つの基本機能", font_size=36, bold=True)

functions = [
    ("知覚", "Perception",
     "テキスト・画像・データなど\n多様な入力を取得・理解する",
     COLOR_ACCENT),
    ("推論", "Reasoning",
     "大規模言語モデル（LLM）が\n文脈を踏まえて思考する",
     COLOR_PURPLE),
    ("計画・実行", "Planning & Action",
     "ツールやAPIを使って\n具体的なアクションを実行する",
     COLOR_GREEN),
    ("学習・記憶", "Memory & Learning",
     "過去のやり取りを保持し\n次の判断に活かす",
     COLOR_ORANGE),
]

for i, (title, en, desc, color) in enumerate(functions):
    x = Inches(0.8 + i * 3.1)
    y = Inches(1.8)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(3.8), RGBColor(0x24, 0x24, 0x3E), 0.05)
    # Color top bar
    add_shape_bg(slide, x, y, Inches(2.8), Pt(6), color, 0)
    # Number
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.4), Inches(1.0), Inches(1.0))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    add_textbox(slide, x + Inches(0.9), y + Inches(0.55), Inches(1.0), Inches(0.7),
                str(i + 1), font_size=32, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.4), Inches(0.5),
                title, font_size=24, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.1), Inches(2.4), Inches(0.3),
                en, font_size=12, alignment=PP_ALIGN.CENTER, color=COLOR_GRAY)
    # Desc
    add_textbox(slide, x + Inches(0.2), y + Inches(2.6), Inches(2.4), Inches(1.0),
                desc, font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

# Key point
key_box = add_shape_bg(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.0), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.15), Inches(9.5), Inches(0.7),
            "RPAは「操作の再生」、AIエージェントは「判断を伴う自律行動」― ここが根本的な違い",
            font_size=17, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 4: GUI操作型 vs API統合型 ― 2つのアプローチ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "AIエージェントの2つのアプローチ", font_size=36, bold=True)

# Left column - GUI型
left_x = Inches(0.8)
card_w = Inches(5.5)
card_l = add_shape_bg(slide, left_x, Inches(1.8), card_w, Inches(4.8), RGBColor(0x24, 0x24, 0x3E), 0.04)
add_shape_bg(slide, left_x, Inches(1.8), card_w, Pt(6), COLOR_ORANGE, 0)
add_textbox(slide, left_x + Inches(0.3), Inches(2.1), card_w - Inches(0.6), Inches(0.5),
            "GUI操作型（画面認識型）", font_size=24, bold=True, color=COLOR_ORANGE)

gui_items = [
    "PC画面のスクリーンショットを撮影して認識",
    "マウス・キーボード操作をシミュレーション",
    "仮想デスクトップ上で動作",
    "既存UIをそのまま利用可能（導入が容易）",
    "画面変更に弱い・動作が比較的遅い",
]
for j, item in enumerate(gui_items):
    add_textbox(slide, left_x + Inches(0.3), Inches(2.9 + j * 0.6), card_w - Inches(0.6), Inches(0.5),
                f"  {item}", font_size=14, color=COLOR_LIGHT_GRAY)

# Right column - API型
right_x = Inches(7.0)
card_r = add_shape_bg(slide, right_x, Inches(1.8), card_w, Inches(4.8), RGBColor(0x24, 0x24, 0x3E), 0.04)
add_shape_bg(slide, right_x, Inches(1.8), card_w, Pt(6), COLOR_GREEN, 0)
add_textbox(slide, right_x + Inches(0.3), Inches(2.1), card_w - Inches(0.6), Inches(0.5),
            "API統合型（システム連携型）", font_size=24, bold=True, color=COLOR_GREEN)

api_items = [
    "ファイルシステムやAPIに直接アクセス",
    "バックグラウンドで高速処理",
    "隔離された仮想環境でセキュアに実行",
    "プラグインや外部サービスとの連携が容易",
    "API非対応のシステムには適用困難",
]
for j, item in enumerate(api_items):
    add_textbox(slide, right_x + Inches(0.3), Inches(2.9 + j * 0.6), card_w - Inches(0.6), Inches(0.5),
                f"  {item}", font_size=14, color=COLOR_LIGHT_GRAY)

# Bottom message
key_box = add_shape_bg(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.85), Inches(9.5), Inches(0.4),
            "企業向けにはAPI統合型が主流になりつつあるが、レガシーシステムではGUI型も有効",
            font_size=15, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: ソフトウェア業界への衝撃 ― SaaSpocalypse
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_RED)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "ソフトウェア業界への衝撃", font_size=36, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "AIエージェントが「ソフトウェアのビジネスモデル」を根本から揺さぶっている",
            font_size=18, color=COLOR_ACCENT2)

impacts = [
    ("UIの価値低下", "AIがAPIで直接操作するなら\nリッチなUIは不要になる",
     "SaaS企業の差別化ポイントが揺らぐ", COLOR_RED),
    ("「席数課金」の崩壊", "AIエージェントは人間ではないが\nソフトウェアを利用する",
     "ユーザー数ベースの価格体系が破綻", COLOR_ORANGE),
    ("統合プラットフォーム化", "個別SaaSを束ねる\nAIオーケストレーターの台頭",
     "単機能SaaSの存在意義が問われる", COLOR_PURPLE),
]

for i, (title, desc, impact, color) in enumerate(impacts):
    y = Inches(2.2 + i * 1.6)
    # Card
    card = add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(1.3), RGBColor(0x24, 0x24, 0x3E), 0.03)
    add_shape_bg(slide, Inches(0.8), y, Pt(6), Inches(1.3), color, 0)
    # Title
    add_textbox(slide, Inches(1.2), y + Inches(0.15), Inches(3.0), Inches(0.5),
                title, font_size=22, bold=True, color=color)
    # Desc
    add_textbox(slide, Inches(4.5), y + Inches(0.15), Inches(4.0), Inches(1.0),
                desc, font_size=14, color=COLOR_LIGHT_GRAY)
    # Impact
    impact_box = add_shape_bg(slide, Inches(8.8), y + Inches(0.2), Inches(3.5), Inches(0.8), RGBColor(0x0A, 0x2A, 0x3E), 0.05)
    add_textbox(slide, Inches(9.0), y + Inches(0.3), Inches(3.2), Inches(0.6),
                impact, font_size=13, color=COLOR_ORANGE, bold=True)

add_textbox(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
            "「SaaSpocalypse」― AIエージェントの台頭でSaaS企業の時価総額が急落する現象が既に起きている",
            font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 6: マルチエージェント ― チームで働くAI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "マルチエージェント ― チームで働くAI", font_size=36, bold=True)

# Central concept
center_box = add_shape_bg(slide, Inches(4.5), Inches(2.0), Inches(4.3), Inches(1.5), COLOR_ACCENT, 0.05)
add_textbox(slide, Inches(4.7), Inches(2.2), Inches(3.9), Inches(0.5),
            "オーケストレーター", font_size=22, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.7), Inches(2.8), Inches(3.9), Inches(0.5),
            "タスクを分割し各エージェントに振り分け", font_size=13, alignment=PP_ALIGN.CENTER)

# Sub agents
agents = [
    ("調査エージェント", "情報収集・\nリサーチ", Inches(0.8), Inches(4.2)),
    ("分析エージェント", "データ分析・\nレポート作成", Inches(4.0), Inches(4.2)),
    ("実行エージェント", "API操作・\nファイル処理", Inches(7.2), Inches(4.2)),
    ("検証エージェント", "品質チェック・\nレビュー", Inches(10.4), Inches(4.2)),
]

for title, desc, x, y in agents:
    agent_box = add_shape_bg(slide, x, y, Inches(2.5), Inches(1.8), RGBColor(0x24, 0x24, 0x3E), 0.05)
    add_shape_bg(slide, x, y, Inches(2.5), Pt(5), COLOR_GREEN, 0)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.3), Inches(2.2), Inches(0.5),
                title, font_size=16, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.2), Inches(0.7),
                desc, font_size=13, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Connecting arrows (simplified as lines)
for ax in [Inches(2.0), Inches(5.2), Inches(8.4), Inches(11.6)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax, Inches(3.55), Inches(0.3), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_ACCENT2
    arrow.line.fill.background()

# Bottom message
key_box = add_shape_bg(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.4), Inches(9.5), Inches(0.6),
            "1つのAIが全部やるのではなく、専門化されたAIがチームとして協働する\n― これが次世代のAIワークフロー",
            font_size=15, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 7: 提言 ― 汎用AIではなく「AIオフィス」を
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "提言：汎用AIツールではなく「AIオフィス」を", font_size=34, bold=True)

# Left side - Problem
prob_box = add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), RGBColor(0x2A, 0x1A, 0x1A), 0.04)
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Pt(5), COLOR_RED, 0)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.5),
            "汎用AIツールの課題", font_size=22, bold=True, color=COLOR_RED)

generic_problems = [
    "毎回プロンプトを考える必要がある",
    "業務知識がないため出力品質にバラつき",
    "社内ルール・文脈を都度説明する手間",
    "使いこなせる人と使えない人の格差",
    "セキュリティ・ガバナンスの管理が困難",
]
for j, item in enumerate(generic_problems):
    add_textbox(slide, Inches(1.2), Inches(2.7 + j * 0.65), Inches(4.8), Inches(0.5),
                f"  {item}", font_size=15, color=COLOR_LIGHT_GRAY)

# Right side - Solution
sol_box = add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), RGBColor(0x0A, 0x2A, 0x1A), 0.04)
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.5), Pt(5), COLOR_GREEN, 0)
add_textbox(slide, Inches(7.4), Inches(1.9), Inches(4.8), Inches(0.5),
            "AIオフィスのあるべき姿", font_size=22, bold=True, color=COLOR_GREEN)

office_features = [
    "業務ごとにプロンプトがプリセット済み",
    "社内ナレッジ・ルールが組み込まれている",
    "誰でも同じ品質のアウトプットを得られる",
    "承認フロー・権限管理と統合されている",
    "業務プロセスに沿ったUIで迷わない",
]
for j, item in enumerate(office_features):
    add_textbox(slide, Inches(7.4), Inches(2.7 + j * 0.65), Inches(4.8), Inches(0.5),
                f"  {item}", font_size=15, color=COLOR_LIGHT_GRAY)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(3.8), Inches(0.6), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLOR_ACCENT2
arrow.line.fill.background()

# Bottom message
add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            "必要なのは「何でもできるAI」ではなく「業務を知っているAI」",
            font_size=16, color=COLOR_ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 8: まとめ ― これからのエンタープライズAI戦略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1.0),
            "まとめ：これからのエンタープライズAI戦略", font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

takeaways = [
    ("RPA → AIエージェント", "操作の再生から\n判断を伴う自律行動へ", COLOR_ACCENT),
    ("API統合が鍵", "UIではなくシステム間連携\nで真の自動化を実現", COLOR_GREEN),
    ("マルチエージェント", "専門化されたAIが\nチームで協働する時代", COLOR_PURPLE),
    ("AIオフィス構想", "汎用ツールではなく\n業務特化型AIを構築", COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.5)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(2.5), RGBColor(0x24, 0x24, 0x3E), 0.04)
    add_shape_bg(slide, x, y, Inches(2.8), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(2.4), Inches(0.6),
                title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(1.0),
                desc, font_size=15, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Final message
msg_box = add_shape_bg(slide, Inches(2.0), Inches(5.5), Inches(9.3), Inches(1.5), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(0.5),
            "AIは道具（ツール）から同僚（コワーカー）へ進化している。",
            font_size=20, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5),
            "しかし企業が求めるのは「万能な同僚」ではなく\n「自社の業務を熟知した専門チーム」である。",
            font_size=17, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# --- Save ---
output_path = "/home/user/harmonic-marketing/ai_agent_evolution.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
