#!/usr/bin/env python3
"""AI時代の承認プロセス - PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xC7)
COLOR_ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent line at top
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

# Title
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            "AI時代の承認プロセス", font_size=48, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.8),
            "人×AIの協働で実現する、速く・正確で・透明な承認ワークフロー",
            font_size=22, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

# Divider
add_accent_line(slide, Inches(5), Inches(4.3), Inches(3.3), COLOR_ACCENT)

# Footer info
add_textbox(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
            "― 経費精算の事例から学ぶ、AI承認の一般フレームワーク ―",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 従来の承認プロセスの課題
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "従来の承認プロセスの課題", font_size=36, bold=True)

problems = [
    ("遅延", "承認者の不在・多忙で\nプロセスが停滞", COLOR_RED),
    ("属人化", "判断基準が承認者ごとに\nばらつき品質が不安定", COLOR_ORANGE),
    ("形骸化", "件数が多く確認が\n形式的になりがち", COLOR_ORANGE),
    ("不透明性", "なぜ承認/却下されたか\n理由が残らない", COLOR_RED),
]

for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    # Card
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(3.5), RGBColor(0x24, 0x24, 0x3E), 0.05)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.4), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, x + Inches(0.9), y + Inches(0.55), Inches(1.0), Inches(0.7),
                "!" if color == COLOR_RED else "?", font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.4), Inches(0.6),
                title, font_size=24, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    # Desc
    add_textbox(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(1.0),
                desc, font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

# Bottom message
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "これらの課題は業種・規模を問わず、あらゆる承認プロセスに共通する構造的問題",
            font_size=15, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: AI承認の基本アーキテクチャ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "AI承認の基本アーキテクチャ", font_size=36, bold=True)

# Flow diagram - horizontal
steps = [
    ("申請", "従業員が\n申請を提出", COLOR_GRAY),
    ("AIチェック", "ルール・AIが\n自動検証", COLOR_ACCENT),
    ("判定", "OK / NG\n自動振り分け", COLOR_PURPLE),
    ("自動承認", "基準適合なら\n即時承認", COLOR_GREEN),
    ("人間レビュー", "要確認案件は\n担当者へ", COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.2)
    # Box
    box = add_shape_bg(slide, x, y, Inches(2.1), Inches(2.2), RGBColor(0x24, 0x24, 0x3E), 0.05)
    # Color top bar
    add_shape_bg(slide, x, y, Inches(2.1), Pt(6), color, 0)
    # Step number
    add_textbox(slide, x + Inches(0.1), y + Inches(0.2), Inches(0.5), Inches(0.4),
                str(i+1), font_size=14, color=color, bold=True)
    # Title
    add_textbox(slide, x + Inches(0.1), y + Inches(0.5), Inches(1.9), Inches(0.5),
                title, font_size=20, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    # Desc
    add_textbox(slide, x + Inches(0.1), y + Inches(1.2), Inches(1.9), Inches(0.8),
                desc, font_size=13, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)
    # Arrow between steps
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.15), y + Inches(0.85), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_ACCENT2
        arrow.line.fill.background()

# Key point box
key_box = add_shape_bg(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(5.15), Inches(9), Inches(0.4),
            "核心：Human-in-the-Loop", font_size=22, bold=True, color=COLOR_ACCENT)
add_textbox(slide, Inches(2.0), Inches(5.65), Inches(9.5), Inches(0.7),
            "AIが処理できる定型判断は自動化し、例外・判断が必要なケースのみ人間にエスカレーション。\n完全自動化ではなく「AIと人の最適な分担」がポイント。",
            font_size=15, color=COLOR_LIGHT_GRAY)


# ============================================================
# Slide 4: AIが担う5つの承認チェック機能
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "AIが担う承認チェック機能（一般化）", font_size=36, bold=True)

checks = [
    ("ルール準拠チェック", "社内規定・法規制への\n適合を自動検証", "規定、ポリシー、\n法令との照合"),
    ("書類・データ完全性", "必要書類の添付漏れ、\n必須項目の記入漏れ検出", "領収書、契約書、\n証明書類の確認"),
    ("整合性・妥当性検証", "金額・日付・分類の\n論理的整合性を確認", "科目、数値範囲、\n日付矛盾の検出"),
    ("コンプライアンス確認", "法的要件・外部基準\nへの準拠を検証", "インボイス番号、\n税務要件の確認"),
    ("異常検知・リスク評価", "過去データとの比較で\n異常値・不正パターンを検出", "統計的外れ値、\n重複申請の検知"),
]

for i, (title, desc, example) in enumerate(checks):
    y = Inches(1.6 + i * 1.1)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_GREEN
    badge.line.fill.background()
    add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.55), Inches(0.45),
                str(i+1), font_size=20, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.6), y, Inches(3.0), Inches(0.6),
                title, font_size=20, bold=True, color=COLOR_ACCENT2)
    # Desc
    add_textbox(slide, Inches(4.8), y, Inches(4.0), Inches(0.7),
                desc, font_size=14, color=COLOR_LIGHT_GRAY)
    # Example
    ex_box = add_shape_bg(slide, Inches(9.2), y, Inches(3.5), Inches(0.65), RGBColor(0x24, 0x24, 0x3E), 0.05)
    add_textbox(slide, Inches(9.4), y + Inches(0.02), Inches(3.2), Inches(0.6),
                example, font_size=12, color=COLOR_GRAY)


# ============================================================
# Slide 5: 導入の3原則
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "AI承認導入の3原則", font_size=36, bold=True)

principles = [
    ("透明性", "Transparency",
     "AIの判断理由を必ず可視化する",
     ["NG理由をコメントとして表示",
      "判断に使用したルールを明示",
      "監査証跡（Audit Trail）を自動記録"]),
    ("段階性", "Gradual Adoption",
     "小さく始めて段階的に拡大する",
     ["まず定型的な承認から自動化",
      "成功実績を積み対象領域を拡大",
      "人間の監視レベルを徐々に調整"]),
    ("協働性", "Human-AI Collaboration",
     "AIと人間の役割を明確に分担する",
     ["AIは定型チェック・一次判断を担当",
      "人間は例外対応・最終判断に集中",
      "全従業員がAIの処理内容を理解"]),
]

for i, (title, en, desc, items) in enumerate(principles):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)
    # Card
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(4.8), RGBColor(0x24, 0x24, 0x3E), 0.04)
    # Number
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
                f"0{i+1}", font_size=14, color=COLOR_PURPLE, bold=True)
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.1), Inches(0.6),
                title, font_size=28, bold=True, color=COLOR_PURPLE)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.4),
                en, font_size=13, color=COLOR_GRAY)
    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.1), Inches(0.6),
                desc, font_size=14, bold=True, color=COLOR_LIGHT_GRAY)
    # Items
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(2.5 + j * 0.6), Inches(3.1), Inches(0.5),
                    f"• {item}", font_size=13, color=COLOR_LIGHT_GRAY)


# ============================================================
# Slide 6: 適用領域の広がり
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "適用領域の広がり ― 経費精算を超えて", font_size=36, bold=True)

domains = [
    ("経費精算", "領収書・規定チェック\n→ 自動承認", "導入済"),
    ("稟議・決裁", "金額・権限に応じた\nルートを自動判定", "拡張可能"),
    ("契約審査", "条項リスク分析\n→ 法務エスカレーション", "拡張可能"),
    ("人事申請", "勤怠・休暇の規定\n適合を自動確認", "拡張可能"),
    ("調達・発注", "予算残・取引先評価\nを自動検証", "拡張可能"),
    ("コード/リリース", "品質基準・テスト結果\nの自動ゲーティング", "拡張可能"),
]

for i, (title, desc, status) in enumerate(domains):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.7)
    # Card
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(2.2), RGBColor(0x24, 0x24, 0x3E), 0.04)
    # Status badge
    badge_color = COLOR_GREEN if status == "導入済" else COLOR_ACCENT
    badge = add_shape_bg(slide, x + Inches(2.2), y + Inches(0.2), Inches(1.3), Inches(0.35), badge_color, 0.5)
    add_textbox(slide, x + Inches(2.2), y + Inches(0.22), Inches(1.3), Inches(0.3),
                status, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.0), Inches(0.5),
                title, font_size=22, bold=True, color=COLOR_ACCENT2)
    # Desc
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(1.0),
                desc, font_size=14, color=COLOR_LIGHT_GRAY)


# ============================================================
# Slide 7: 導入ロードマップ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "導入ロードマップ", font_size=36, bold=True)

phases = [
    ("Phase 1", "1〜2ヶ月", "アセスメント",
     ["現行承認フローの可視化", "自動化対象の選定", "ルール・規定の整理"]),
    ("Phase 2", "2〜3ヶ月", "PoC・パイロット",
     ["小規模部門で試行", "AIモデルのチューニング", "ユーザーフィードバック収集"]),
    ("Phase 3", "1〜2ヶ月", "本番展開",
     ["全社展開・教育", "モニタリング体制構築", "継続的改善サイクル確立"]),
]

# Timeline bar
bar = add_shape_bg(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(0.15), COLOR_ACCENT, 0.5)

for i, (phase, period, title, items) in enumerate(phases):
    x = Inches(1.5 + i * 3.8)
    # Circle on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4), Inches(2.0), Inches(0.5), Inches(0.5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLOR_ACCENT
    dot.line.fill.background()
    add_textbox(slide, x + Inches(1.4), Inches(2.05), Inches(0.5), Inches(0.4),
                str(i+1), font_size=18, bold=True, alignment=PP_ALIGN.CENTER)
    # Card below
    card = add_shape_bg(slide, x, Inches(3.0), Inches(3.4), Inches(3.8), RGBColor(0x24, 0x24, 0x3E), 0.04)
    add_textbox(slide, x + Inches(0.3), Inches(3.2), Inches(2.8), Inches(0.4),
                phase, font_size=14, color=COLOR_ACCENT, bold=True)
    add_textbox(slide, x + Inches(0.3), Inches(3.6), Inches(2.8), Inches(0.5),
                title, font_size=24, bold=True, color=COLOR_WHITE)
    add_textbox(slide, x + Inches(1.7), Inches(3.25), Inches(1.5), Inches(0.3),
                period, font_size=12, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), Inches(4.3 + j * 0.55), Inches(2.8), Inches(0.5),
                    f"• {item}", font_size=14, color=COLOR_LIGHT_GRAY)


# ============================================================
# Slide 8: まとめ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1.0),
            "まとめ：AI承認がもたらす変革", font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

benefits = [
    ("スピード", "承認リードタイムを\n最大90%短縮", COLOR_ACCENT),
    ("正確性", "ヒューマンエラーを\n大幅に削減", COLOR_GREEN),
    ("透明性", "全判断の根拠が\n記録・追跡可能", COLOR_PURPLE),
    ("スケーラビリティ", "件数増加にも\nコスト比例せず対応", COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.0)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(2.2), RGBColor(0x24, 0x24, 0x3E), 0.04)
    # Color accent top
    add_shape_bg(slide, x, y, Inches(2.8), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(2.4), Inches(0.5),
                title, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.8),
                desc, font_size=15, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Final message
msg_box = add_shape_bg(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.0), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.8), Inches(5.85), Inches(7.7), Inches(0.7),
            "AIは承認者を置き換えるのではなく、承認者がより価値の高い判断に集中できる環境をつくる",
            font_size=18, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# --- Save ---
output_path = "/home/user/harmonic-marketing/ai_approval_process.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
