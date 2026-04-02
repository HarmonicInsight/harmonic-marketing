#!/usr/bin/env python3
"""建設業法・入契法改正（令和6年法律第49号）解説スライド - PPTX生成スクリプト

国土交通省サイトの資料画像を活用した建設業向け解説プレゼンテーション。
出典: https://www.mlit.go.jp/tochi_fudousan_kensetsugyo/const/tochi_fudousan_kensetsugyo_const_tk1_000001_00033.html
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette (建設業・官公庁向けの落ち着いた配色) ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x6B, 0xB6)  # 国交省ブルー
COLOR_ACCENT2 = RGBColor(0x48, 0xA0, 0xD0)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xE8, 0x8C, 0x00)
COLOR_RED = RGBColor(0xDC, 0x38, 0x38)
COLOR_YELLOW = RGBColor(0xF5, 0xC5, 0x18)
COLOR_CARD = RGBColor(0x24, 0x24, 0x3E)
COLOR_HIGHLIGHT_BG = RGBColor(0x0A, 0x20, 0x3E)

ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets", "construction_law")


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo",
                   bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_notes(slide, text):
    """スライドにノート（ナレーション台本）を追加"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_image_slide(slide, image_filename, title_text, subtitle_text=None,
                    accent_color=COLOR_ACCENT):
    """サイト画像をアスペクト比を保持して中央配置するスライドを生成"""
    add_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, accent_color)

    # タイトル
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
                title_text, font_size=32, bold=True, color=COLOR_WHITE)

    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.5),
                    subtitle_text, font_size=16, color=COLOR_ACCENT2)

    # 画像をアスペクト比を保持して配置
    img_path = os.path.join(ASSETS_DIR, image_filename)
    from PIL import Image as PILImage
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h  # 全画像 ≈ 1.445

    img_top = Inches(1.55) if subtitle_text else Inches(1.3)
    max_height = Inches(5.5) if subtitle_text else Inches(5.7)
    max_width = Inches(11.5)

    # アスペクト比を保持してフィット
    if max_width / max_height > aspect:
        # 高さ基準
        final_height = max_height
        final_width = int(max_height * aspect)
    else:
        # 幅基準
        final_width = max_width
        final_height = int(max_width / aspect)

    # 水平方向の中央揃え
    slide_w = prs.slide_width
    img_left = int((slide_w - final_width) / 2)

    # 画像の白背景用カード（画像より少し大きく）
    pad = Inches(0.1)
    add_shape_bg(slide, img_left - pad, img_top - pad,
                 final_width + pad * 2, final_height + pad * 2,
                 COLOR_WHITE, 0.02)

    slide.shapes.add_picture(img_path, img_left, img_top,
                             final_width, final_height)

    # 出典表示
    add_textbox(slide, Inches(7.5), Inches(7.1), Inches(5.5), Inches(0.3),
                "出典：国土交通省", font_size=10, color=COLOR_GRAY,
                alignment=PP_ALIGN.RIGHT)


# ============================================================
# Slide 1: タイトルスライド
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

# 上部ラベル
label_bg = add_shape_bg(slide, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.5),
                        COLOR_ACCENT, 0.5)
add_textbox(slide, Inches(4.5), Inches(1.22), Inches(4.3), Inches(0.45),
            "建設業者の皆さまへ", font_size=18, bold=True,
            alignment=PP_ALIGN.CENTER)

# メインタイトル
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
            "建設業法・入契法 改正のポイント", font_size=52, bold=True,
            alignment=PP_ALIGN.CENTER)

# サブタイトル
add_textbox(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(0.7),
            "令和6年法律第49号（令和6年6月14日公布）",
            font_size=24, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

# 区切り線
add_accent_line(slide, Inches(5), Inches(4.8), Inches(3.3), COLOR_ACCENT)

# 施行スケジュール概要
schedule_box = add_shape_bg(slide, Inches(2.5), Inches(5.3), Inches(8.3), Inches(1.5),
                            COLOR_HIGHLIGHT_BG, 0.03)

schedules = [
    ("第1段階", "令和6年9月1日施行", "公布から3ヶ月"),
    ("第2段階", "令和6年12月13日施行", "公布から6ヶ月"),
    ("第3段階", "公布から1年6ヶ月以内", "施行日未定"),
]

for i, (phase, date, note) in enumerate(schedules):
    x = Inches(2.8 + i * 2.7)
    y = Inches(5.45)
    # フェーズ名
    add_textbox(slide, x, y, Inches(2.3), Inches(0.35),
                phase, font_size=13, color=COLOR_ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    # 日付
    add_textbox(slide, x, y + Inches(0.35), Inches(2.3), Inches(0.4),
                date, font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
    # 備考
    add_textbox(slide, x, y + Inches(0.8), Inches(2.3), Inches(0.3),
                note, font_size=11, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide, """皆さん、こんにちは。
今回は、建設業法と入契法（にゅうけいほう）の改正について、ポイントをわかりやすくお話しします。
令和6年の6月に公布されたこの改正、実は建設業で働く皆さんにとって、かなり大きな変化をもたらすものです。
施行は3段階に分かれていて、一番早いものは令和6年9月にもうスタートしています。
「何が変わるの？」「うちの会社は何をすればいいの？」――そんな疑問に答えていきますので、ぜひ最後までご覧ください。""")


# ============================================================
# Slide 2: 改正の背景 - 建設業の現状
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_overview.jpg",
                "建設業をめぐる現状と改正の背景",
                "建設投資・許可業者数・就業者数の推移")

add_notes(slide, """では、そもそもなぜこの改正が必要だったのか。
このグラフを見てください。建設業の現状がひと目でわかります。
まず、働く人がどんどん減っています。ピークから約3割減。これが一番深刻ですよね。
業者の数も2割ほど減っていますし、投資額もピーク時には遠く及びません。
その一方で、資材の値段はここ数年で急激に上がっている。
つまり、「人は減る、モノは高くなる、でも仕事はある」という状態なんですね。
このままだと、しわ寄せが現場の労務費に来てしまう。
この危機感から、今回の法改正が動き出しました。""")


# ============================================================
# Slide 3: 改正の3つの柱（概要）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "改正の3つの柱", font_size=36, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "建設業の持続的発展に向けた制度的基盤の整備",
            font_size=16, color=COLOR_ACCENT2)

pillars = [
    ("1", "処遇改善",
     "労働者の処遇確保",
     [
         "処遇確保の努力義務化",
         "中央建設業審議会が「労務費の基準」を作成",
         "著しく低い労務費での見積を禁止",
         "不当に低い請負代金の禁止を受注者にも拡大",
     ],
     COLOR_GREEN),
    ("2", "資材高騰対策",
     "労務費へのしわ寄せ防止",
     [
         "請負代金の変更方法を契約書に明記",
         "資材高騰の「おそれ情報」通知を義務化",
         "変更協議への誠実対応を努力義務化",
         "公共発注者には協議に応じる義務",
     ],
     COLOR_ORANGE),
    ("3", "働き方改革・生産性向上",
     "担い手の確保・育成",
     [
         "著しく短い工期での契約を禁止",
         "工期変更の協議円滑化",
         "ICT活用で専任義務を合理化",
         "現場技術者の専任義務の見直し",
     ],
     COLOR_ACCENT),
]

for i, (num, title, subtitle, items, color) in enumerate(pillars):
    x = Inches(0.6 + i * 4.15)
    y = Inches(1.9)

    # カード
    card = add_shape_bg(slide, x, y, Inches(3.9), Inches(5.0), COLOR_CARD, 0.03)
    # 上部カラーバー
    add_shape_bg(slide, x, y, Inches(3.9), Pt(6), color, 0)

    # 番号バッジ
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.4),
                                   Inches(0.6), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_textbox(slide, x + Inches(0.3), y + Inches(0.45), Inches(0.6), Inches(0.5),
                num, font_size=24, bold=True, alignment=PP_ALIGN.CENTER)

    # タイトル
    add_textbox(slide, x + Inches(1.1), y + Inches(0.4), Inches(2.5), Inches(0.5),
                title, font_size=26, bold=True, color=color)
    add_textbox(slide, x + Inches(1.1), y + Inches(0.9), Inches(2.5), Inches(0.4),
                subtitle, font_size=13, color=COLOR_GRAY)

    # 項目リスト
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(1.6 + j * 0.75),
                    Inches(3.3), Inches(0.7),
                    f"  {item}", font_size=14, color=COLOR_LIGHT_GRAY)
        # 先頭のドット
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.35), y + Inches(1.72 + j * 0.75),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

add_notes(slide, """今回の改正は、大きく3つの柱で構成されています。
1つ目が「処遇改善」。ざっくり言うと、現場で働く人にちゃんとお金が届くようにする仕組みです。
2つ目が「資材高騰対策」。鉄やコンクリートの値段が上がった分を、きちんと価格に転嫁できるようにするルールです。
3つ目が「働き方改革と生産性向上」。無理な工期をなくして、ICTも活用しようという話ですね。
どれも建設業の現場に直結する内容ですので、それぞれ詳しく見ていきましょう。""")


# ============================================================
# Slide 4: 処遇改善の詳細
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_treatment.jpg",
                "柱1：処遇改善",
                "建設業者の責務・労務費の確保と行き渡り・不当に低い請負代金の禁止",
                accent_color=COLOR_GREEN)

add_notes(slide, """まず1つ目の柱、「処遇改善」です。
ポイントは3つあります。
1つ目、建設業者に対して「労働者の処遇を確保すること」が努力義務になりました。
これまでも当然やるべきことでしたが、法律でちゃんと位置づけられたのは大きいですね。
2つ目、中央建設業審議会が「労務費の基準」を作ります。
これによって「著しく低い労務費」での見積りの依頼や提出が禁止されます。
要するに、「安くしてよ」と無理なお願いをすること自体がNGになるわけです。
3つ目、不当に低い請負代金の禁止が、発注者だけでなく受注者にも広がりました。
下請けに丸投げで安く叩く、というのも許されなくなります。""")


# ============================================================
# Slide 5: 労務費の基準
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_labor_cost.jpg",
                "「著しく低い労務費等」と「不当に低い請負代金」の禁止",
                "中央建設業審議会が「労務費に関する基準」を作成・勧告",
                accent_color=COLOR_GREEN)

add_notes(slide, """こちらがその仕組みの詳細です。
注文者、つまり発注側が「著しく低い労務費」で見積りを依頼すると、これは法律違反になります。
同じように、受注者が自ら低すぎる労務費で見積りを出すことも禁止です。
じゃあ「著しく低い」って何が基準になるの？という話ですが、
中央建設業審議会が「労務費に関する基準」を作成して公表する仕組みになっています。
違反した場合は、発注者には勧告・公表、受注者には指導・監督が行われます。
つまり、元請けも下請けも、お互いに適正な価格を守らないといけない。
「安ければいい」という時代は終わった、ということですね。""")


# ============================================================
# Slide 6: 資材高騰対策
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_material_price.jpg",
                "柱2：資材高騰に伴う労務費のしわ寄せ防止",
                "契約前後のルールを整備し、価格転嫁協議の円滑化を実現",
                accent_color=COLOR_ORANGE)

add_notes(slide, """2つ目の柱、「資材高騰への対策」です。
ここ数年、鉄鋼や木材など、資材の価格が急上昇していますよね。
問題は、その値上がり分がどこに行くかということです。
契約金額が変わらないまま資材だけ高くなると、結局、労務費が削られてしまう。
これを防ぐために、契約のルールが2つ整備されました。
まず「契約前のルール」。請負代金の変更方法を、契約書に必ず書くことが義務化されました。
さらに、受注者は資材が高騰しそうだという「おそれ情報」を、事前に発注者に通知する義務があります。
次に「契約後のルール」。通知を受けた発注者は、価格変更の協議に誠実に対応する努力義務を負います。
公共工事の発注者には、協議に応じる義務が課されています。
つまり、「値上がりしたけど言い出せない」という状況をなくす仕組みですね。""")


# ============================================================
# Slide 7: おそれ情報の通知制度
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_risk_info.jpg",
                "価格転嫁協議の円滑化ルールの詳細",
                "「おそれ情報」の通知制度と変更協議のプロセス",
                accent_color=COLOR_ORANGE)

add_notes(slide, """こちらが具体的なルールの流れです。
契約前の段階では、受注者が「資材高騰のおそれがある」と判断したら、注文者に通知します。
このとき大事なのが、「根拠情報」をちゃんと添えること。
たとえば、業界紙の記事や、メーカーの価格改定のお知らせ、統計データなどですね。
「なんとなく上がりそう」ではなく、客観的なデータで示す必要があります。
そして実際に資材が高騰したら、契約書に書いた変更方法に従って、価格変更の協議に入ります。
注文者は誠実に協議に応じなければなりません。
この通知や協議の記録は、書面やメールで双方が保存しておくことが求められます。
後から「言った、言わない」にならないようにする仕組みですね。""")


# ============================================================
# Slide 8: 働き方改革・生産性向上
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_workstyle.jpg",
                "柱3：働き方改革と生産性向上",
                "工期ダンピング対策の強化とICTを活用した現場管理の効率化",
                accent_color=COLOR_ACCENT)

add_notes(slide, """3つ目の柱は「働き方改革と生産性向上」。ここは大きく2つに分かれます。
まず「働き方改革」のほうから。
目玉は、「著しく短い工期」での契約が禁止されたことです。
いわゆる「工期ダンピング」ですね。無理な工期で受けて、現場に長時間労働を強いる。
これが法律ではっきりNGになりました。
新たに注文者にも工期に関する責任が課され、違反した建設業者には指導・監督が入ります。
また、資材が手に入らないなどの事情で工期の変更が必要なときは、協議をスムーズに行う制度も導入されました。
次に「生産性向上」のほう。
ICT、つまりタブレットやウェアラブルカメラを使った現場管理の効率化を後押しする内容です。
これについては次のスライドで詳しくお話しします。""")


# ============================================================
# Slide 9: ICT活用による専任義務の合理化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "slide_productivity.jpg",
                "ICT活用による監理技術者等の専任義務の合理化",
                "タブレット・ウェアラブルカメラ等の活用で遠隔施工管理を実現",
                accent_color=COLOR_ACCENT)

add_notes(slide, """これが注目のICT活用の話です。
今まで、監理技術者や主任技術者は、基本的に現場に「専任」で張りついている必要がありました。
でも、タブレットで図面や写真を共有したり、ウェアラブルカメラでリアルタイムに映像を送れる時代ですよね。
そこで今回の改正では、ICTを活用することを条件に、この専任義務が合理化されました。
具体的には、ICT機器を使って適切に施工管理ができる体制が整っていれば、
1人の技術者が複数の現場を兼務できるようになります。
人手不足が深刻な中、これは現場にとって非常にありがたい変更ではないでしょうか。
ただし、条件や要件がありますので、導入を検討される場合は国交省の指針をしっかり確認してください。""")


# ============================================================
# Slide 10: まとめ - 建設業者が対応すべきポイント
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.5), Inches(10.3), Inches(0.9),
            "建設業者が対応すべきポイント", font_size=40, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(0.5),
            "改正法への対応チェックリスト", font_size=18, color=COLOR_ACCENT2,
            alignment=PP_ALIGN.CENTER)

actions = [
    ("処遇改善への対応", "労務費の基準を確認し、\n適正な労務費の確保を", COLOR_GREEN,
     ["労務費の基準を把握・遵守", "処遇確保の取組を実施", "著しく低い労務費での見積を回避"]),
    ("契約書の見直し", "変更方法の明記と\nおそれ情報への対応", COLOR_ORANGE,
     ["請負代金の変更方法を契約書に明記", "資材高騰のおそれ情報の通知体制整備", "根拠資料の収集・保存の仕組み構築"]),
    ("工期の適正化", "適正工期の確保と\n週休2日の推進", COLOR_RED,
     ["著しく短い工期での受注を回避", "長時間労働の是正に取り組む", "工期変更時の協議記録を残す"]),
    ("ICT活用の検討", "デジタル技術で\n生産性向上を", COLOR_ACCENT,
     ["タブレット・カメラ等のICT導入検討", "専任義務の合理化条件の確認", "遠隔施工管理体制の構築"]),
]

for i, (title, desc, color, items) in enumerate(actions):
    x = Inches(0.5 + i * 3.15)
    y = Inches(2.2)

    # カード
    card = add_shape_bg(slide, x, y, Inches(2.95), Inches(4.8), COLOR_CARD, 0.03)
    # カラーバー
    add_shape_bg(slide, x, y, Inches(2.95), Pt(5), color, 0)

    # タイトル
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.55), Inches(0.5),
                title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)

    # 説明
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.55), Inches(0.8),
                desc, font_size=13, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # チェック項目
    for j, item in enumerate(items):
        check_y = y + Inches(1.9 + j * 0.85)
        # チェックボックス風
        check_box = add_shape_bg(slide, x + Inches(0.25), check_y + Inches(0.05),
                                 Inches(0.25), Inches(0.25), COLOR_HIGHLIGHT_BG, 0.1)
        add_textbox(slide, x + Inches(0.6), check_y, Inches(2.1), Inches(0.7),
                    item, font_size=12, color=COLOR_LIGHT_GRAY)

# フッターメッセージ
footer_box = add_shape_bg(slide, Inches(2.0), Inches(7.05), Inches(9.3), Inches(0.35),
                          COLOR_HIGHLIGHT_BG, 0.5)
add_textbox(slide, Inches(2.0), Inches(7.07), Inches(9.3), Inches(0.3),
            "出典：国土交通省「建設業法・入契法改正（令和6年法律第49号）について」",
            font_size=11, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide, """最後に、建設業者の皆さんが具体的に何をすればいいのか、4つのポイントにまとめました。
1つ目、処遇改善への対応。労務費の基準が公表されたら、まずそれを確認してください。
自社の見積りが基準を下回っていないか、しっかりチェックしましょう。
2つ目、契約書の見直し。請負代金の変更方法が契約書に入っているか、今一度確認してください。
おそれ情報の通知体制も整えておくと安心です。根拠資料の収集・保存の仕組みも作っておきましょう。
3つ目、工期の適正化。無理な工期での受注は、もう法律で禁止されています。
週休2日の確保や長時間労働の是正にも、積極的に取り組んでいきましょう。
4つ目、ICT活用の検討。タブレットやカメラの導入で、技術者の専任義務が緩和される可能性があります。
人手不足への対策としても、ぜひ検討してみてください。
以上が、建設業法・入契法改正のポイントでした。
この改正は、建設業で働く皆さんの処遇を守り、持続可能な業界にしていくための大きな一歩です。
ぜひ、自社の対応を今のうちから進めていきましょう。ご視聴ありがとうございました。""")


# --- Save ---
output_path = os.path.join(os.path.dirname(__file__), "construction_law_reform.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
