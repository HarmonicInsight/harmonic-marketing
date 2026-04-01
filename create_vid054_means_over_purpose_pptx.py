#!/usr/bin/env python3
"""VID-054: DXが失敗する会社の共通点 - PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xC7)
COLOR_ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
COLOR_CARD = RGBColor(0x24, 0x24, 0x3E)
COLOR_HIGHLIGHT_BG = RGBColor(0x0A, 0x2A, 0x3E)


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: タイトル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_RED)

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            "DXが失敗する会社の共通点", font_size=48, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.8),
            "ツール導入は目的ではない──建設業DXで最初に倒すべき壁",
            font_size=22, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.3), Inches(3.3), COLOR_RED)

add_textbox(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
            "― 建設業DXの現実シリーズ ―",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: よくある失敗パターン
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_RED)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "ツールは変わった。でも業務は変わっていない。", font_size=36, bold=True)

failures = [
    ("AI議事録ツール", "導入した", "誰も見返さない", "📝"),
    ("クラウド工程管理", "導入した", "Excelと二重管理", "📊"),
    ("チャットツール", "導入した", "電話・LINEも併用\n情報が三重化", "💬"),
]

for i, (tool, action, result, icon) in enumerate(failures):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)

    # Card
    add_shape_bg(slide, x, y, Inches(3.7), Inches(4.0), COLOR_CARD, 0.04)

    # Tool name
    add_textbox(slide, x + Inches(0.3), y + Inches(0.4), Inches(3.1), Inches(0.6),
                tool, font_size=24, bold=True, color=COLOR_ACCENT2,
                alignment=PP_ALIGN.CENTER)

    # Action
    add_textbox(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.5),
                action, font_size=16, color=COLOR_GREEN,
                alignment=PP_ALIGN.CENTER)

    # Arrow down
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   x + Inches(1.5), y + Inches(1.7),
                                   Inches(0.7), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_RED
    arrow.line.fill.background()

    # Result
    result_bg = add_shape_bg(slide, x + Inches(0.3), y + Inches(2.4),
                             Inches(3.1), Inches(1.2), RGBColor(0x3E, 0x1A, 0x1A), 0.04)
    add_textbox(slide, x + Inches(0.4), y + Inches(2.55), Inches(2.9), Inches(0.9),
                result, font_size=18, bold=True, color=COLOR_RED,
                alignment=PP_ALIGN.CENTER)

# Bottom message
msg = add_shape_bg(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.6),
            "ツールは変わった。でも業務は何一つ変わっていない。",
            font_size=18, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: 原因──業務が整理されていない
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "原因：業務の手順が整理されていない", font_size=36, bold=True)

# Left side - business flow
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
            "業務は「一定の手順」で動いている", font_size=20, bold=True,
            color=COLOR_ACCENT2)

flow_steps = ["見積作成", "承認取得", "契約", "着工", "検収", "請求"]
for i, step in enumerate(flow_steps):
    x = Inches(0.8 + i * 1.8)
    y = Inches(2.3)
    box = add_shape_bg(slide, x, y, Inches(1.5), Inches(0.7), COLOR_CARD, 0.05)
    add_textbox(slide, x, y + Inches(0.1), Inches(1.5), Inches(0.5),
                step, font_size=14, bold=True, alignment=PP_ALIGN.CENTER,
                color=COLOR_WHITE)
    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(1.52), y + Inches(0.15),
                                       Inches(0.25), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_ACCENT2
        arrow.line.fill.background()

# Right side - problems
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(0.5),
            "しかし現実は……", font_size=20, bold=True, color=COLOR_ORANGE)

problems = [
    "Aさんのやり方とBさんのやり方が違う",
    "現場ごと、時期ごとにバラバラ",
    "手順が可視化・標準化されていない",
    "ツールに乗せるべき「業務」が定義されていない",
]

for i, prob in enumerate(problems):
    y = Inches(4.2 + i * 0.6)
    # X mark
    add_textbox(slide, Inches(0.8), y, Inches(0.5), Inches(0.5),
                "✕", font_size=18, bold=True, color=COLOR_RED)
    add_textbox(slide, Inches(1.4), y, Inches(10), Inches(0.5),
                prob, font_size=16, color=COLOR_LIGHT_GRAY)

# Key message
msg = add_shape_bg(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(6.4), Inches(9.3), Inches(0.6),
            "業務が定義されていない状態で、どんなツールを入れても定着しない",
            font_size=18, bold=True, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 4: 手段の目的化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            '典型的な「手段の目的化」', font_size=36, bold=True)

# Table-like comparison
headers = ["", "本来あるべき姿", "実際に起きていること"]
header_colors = [COLOR_GRAY, COLOR_GREEN, COLOR_RED]

for i, (header, hcolor) in enumerate(zip(headers, header_colors)):
    x = Inches(0.8 + i * 4.1)
    add_textbox(slide, x, Inches(1.8), Inches(3.7), Inches(0.5),
                header, font_size=16, bold=True, color=hcolor,
                alignment=PP_ALIGN.CENTER)

rows = [
    ("目的", "業務を効率化・\n標準化したい", "ツール導入自体が\nゴールになっている"),
    ("手段", "そのために\nツールを活用する", "業務整理は\n「面倒だから後で」"),
    ("結果", "業務が変わり\n生産性が上がる", "ツールが増えただけで\n何も変わらない"),
]

for i, (label, ideal, reality) in enumerate(rows):
    y = Inches(2.5 + i * 1.5)

    # Label
    label_box = add_shape_bg(slide, Inches(0.8), y, Inches(3.7), Inches(1.2),
                             COLOR_CARD, 0.04)
    add_textbox(slide, Inches(0.8), y + Inches(0.25), Inches(3.7), Inches(0.7),
                label, font_size=22, bold=True, color=COLOR_WHITE,
                alignment=PP_ALIGN.CENTER)

    # Ideal
    ideal_box = add_shape_bg(slide, Inches(4.9), y, Inches(3.7), Inches(1.2),
                             RGBColor(0x0A, 0x2E, 0x1A), 0.04)
    add_textbox(slide, Inches(5.1), y + Inches(0.15), Inches(3.3), Inches(0.9),
                ideal, font_size=16, color=COLOR_GREEN,
                alignment=PP_ALIGN.CENTER)

    # Reality
    reality_box = add_shape_bg(slide, Inches(9.0), y, Inches(3.7), Inches(1.2),
                               RGBColor(0x3E, 0x1A, 0x1A), 0.04)
    add_textbox(slide, Inches(9.2), y + Inches(0.15), Inches(3.3), Inches(0.9),
                reality, font_size=16, color=COLOR_RED,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: なぜ業務整理をやらないのか
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "なぜ業務整理をやらないのか", font_size=36, bold=True)

reasons = [
    ("① 面倒だから", "地味で時間がかかる作業\n現場ヒアリング・手順書き出し\n例外の洗い出し・関係者合意",
     COLOR_ORANGE),
    ("② リーダーがいない", "部門横断の取り組みが必要\nしかし旗振り役がいない\n結果「見えやすい」ツール導入だけ進む",
     COLOR_RED),
    ("③ コンサルも\n  業務整理をしない", "見るべき指標は業界で決まっている\nしかし体系化できるコンサルが少ない\n「お客様が見たいもの」を見せるだけ",
     COLOR_PURPLE),
]

for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)

    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(4.5), COLOR_CARD, 0.04)
    add_shape_bg(slide, x, y, Inches(3.7), Pt(6), color, 0)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(1.35), y + Inches(0.4),
                                    Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, x + Inches(1.35), y + Inches(0.55), Inches(1.0), Inches(0.7),
                str(i + 1), font_size=36, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.3), y + Inches(1.6), Inches(3.1), Inches(0.8),
                title, font_size=20, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)

    add_multi_text(slide, x + Inches(0.3), y + Inches(2.5), Inches(3.1), Inches(1.8),
                   desc.split("\n"), font_size=14, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.6)

# Bottom callout
msg = add_shape_bg(slide, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.6),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(6.65), Inches(9.3), Inches(0.5),
            "見るべきものを示すのがコンサルタントの仕事。御用聞きはコンサルティングではない。",
            font_size=16, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 6: 正しい順序
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DXで成果を出す正しい順序", font_size=36, bold=True)

steps = [
    ("01", "業務整理", "手順を可視化\n標準化する", COLOR_ACCENT),
    ("02", "情報統合", "Excel・紙・口頭の\n情報を一箇所に集める", COLOR_ACCENT2),
    ("03", "ツール選定", "整理された業務に\n合うツールを選ぶ", COLOR_GREEN),
    ("04", "AI活用", "統合データの上で\nAIを活用する", COLOR_PURPLE),
]

# Timeline bar
bar = add_shape_bg(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.15),
                   COLOR_ACCENT, 0.5)

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(1.0 + i * 2.9)

    # Circle on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + Inches(0.85), Inches(2.25),
                                 Inches(0.6), Inches(0.6))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_textbox(slide, x + Inches(0.85), Inches(2.3), Inches(0.6), Inches(0.5),
                num, font_size=16, bold=True, alignment=PP_ALIGN.CENTER)

    # Card below
    card = add_shape_bg(slide, x, Inches(3.2), Inches(2.6), Inches(2.6),
                        COLOR_CARD, 0.04)
    add_shape_bg(slide, x, Inches(3.2), Inches(2.6), Pt(6), color, 0)

    add_textbox(slide, x + Inches(0.2), Inches(3.5), Inches(2.2), Inches(0.5),
                title, font_size=22, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.2), Inches(4.2), Inches(2.2), Inches(1.2),
                   desc.split("\n"), font_size=14, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.5)

# Warning box
warn = add_shape_bg(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8),
                    RGBColor(0x3E, 0x1A, 0x1A), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.25), Inches(2.0), Inches(0.6),
            "⚠ よくある失敗", font_size=16, bold=True, color=COLOR_RED)
add_textbox(slide, Inches(4.0), Inches(6.35), Inches(7.5), Inches(0.5),
            "ステップ①②を飛ばして③④から始める → 確実に失敗する",
            font_size=16, bold=True, color=COLOR_ORANGE)


# ============================================================
# Slide 7: まとめ・CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1.0),
            "まとめ：DXの第一歩はツール選びではない",
            font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

# 3 key takeaways
takeaways = [
    ("業務整理が先", "ツールを入れる前に\n手順を可視化・標準化", COLOR_ACCENT),
    ("手段の目的化に\n気づく", "ツール導入がゴールに\nなっていないか確認", COLOR_ORANGE),
    ("正しい順序を守る", "業務整理→情報統合\n→ツール→AI", COLOR_GREEN),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.5)
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(2.2), COLOR_CARD, 0.04)
    add_shape_bg(slide, x, y, Inches(3.7), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.4), Inches(3.1), Inches(0.7),
                title, font_size=22, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.8),
                   desc.split("\n"), font_size=15, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.4)

# CTA box
cta = add_shape_bg(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.8),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(5.35), Inches(9), Inches(0.4),
            "HARMONIC insight ── 建設業の業務整理からDX実現まで",
            font_size=20, bold=True, color=COLOR_ACCENT)

cta_items = [
    "自分で始めたい方 → note記事「DXが失敗する会社の共通点」",
    "体系的に学びたい方 → SIPOフレームワーク実践ワークショップ",
    "支援が必要な方 → お問い合わせ（概要欄リンク）",
]
for i, item in enumerate(cta_items):
    add_textbox(slide, Inches(2.5), Inches(5.85 + i * 0.45), Inches(8.5), Inches(0.4),
                f"▸ {item}", font_size=14, color=COLOR_LIGHT_GRAY)


# --- Save ---
output_path = "/home/user/harmonic-marketing/content/youtube/slides/VID-054_means_over_purpose.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
