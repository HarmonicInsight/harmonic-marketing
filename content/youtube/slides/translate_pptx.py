"""
Generic PPTX translator: copies a .pptx, replaces Japanese text with English.
Usage: import and call translate_file(src, dst, translations)
"""
import shutil
from pptx import Presentation
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


def translate_file(src: str, dst: str, translations: dict) -> dict:
    """
    Copy src to dst and replace text according to translations dict.
    Keys in translations can be plain text or contain \x0b for line breaks.
    Returns {'replaced': int, 'missed': list[str]}
    """
    shutil.copy2(src, dst)
    prs = Presentation(dst)
    replaced = 0
    missed = []

    def process_paragraph(para):
        nonlocal replaced
        full = para.text
        if not full.strip():
            return

        # Try exact match first
        key = full.strip()
        if key in translations:
            en = translations[key]
        elif full in translations:
            en = translations[full]
        else:
            # Try normalized (collapse whitespace)
            norm = ' '.join(full.split())
            found = False
            for k, v in translations.items():
                if ' '.join(k.split()) == norm:
                    en = v
                    found = True
                    break
            if not found:
                if any(ord(c) > 0x3000 for c in full):
                    missed.append(full.strip())
                return

        # Check if paragraph has line breaks (\x0b)
        has_breaks = '\x0b' in full

        if has_breaks:
            en_parts = en.split('\x0b') if '\x0b' in en else [en]
            elements = list(para._p)
            part_idx = 0
            first_run = True

            for elem in elements:
                tag = etree.QName(elem.tag).localname if isinstance(elem.tag, str) else ''
                if tag == 'r':
                    t_elem = elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    if t_elem is not None:
                        if first_run and part_idx < len(en_parts):
                            t_elem.text = en_parts[part_idx]
                            first_run = False
                        else:
                            t_elem.text = ''
                elif tag == 'br':
                    part_idx += 1
                    first_run = True
        else:
            # Simple case: put translation in first run, clear rest
            if para.runs:
                para.runs[0].text = en
                for r in para.runs[1:]:
                    r.text = ''
        replaced += 1

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    process_paragraph(para)
        # Also translate notes
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                process_paragraph(para)

    prs.save(dst)
    return {'replaced': replaced, 'missed': list(set(missed))}
