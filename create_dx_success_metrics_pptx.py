#!/usr/bin/env python3
"""DX成功指標 ― 受験とダイエットに学ぶDX運用の本質 - PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xC7)
COLOR_ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
COLOR_PINK = RGBColor(0xEC, 0x48, 0x99)


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: タイトル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "DX成功の方程式", font_size=52, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
            "受験とダイエットに学ぶ、DXツール導入後の「本当の勝負」",
            font_size=24, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.2), Inches(3.3), COLOR_ACCENT)

add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.5),
            "― ツールを買っただけでは、合格も減量もできない ―",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: DXの現状 ― よくある失敗パターン
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_RED)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DXの現状 ― よくある失敗パターン", font_size=36, bold=True)

failures = [
    ("導入して満足", "ツールを導入した\n時点でゴールだと\n錯覚してしまう", COLOR_RED),
    ("運用が定着しない", "現場が使わない、\n旧来のやり方に\n戻ってしまう", COLOR_ORANGE),
    ("効果が測れない", "KPIが曖昧で\n成果を証明\nできない", COLOR_ORANGE),
    ("経営に響かない", "結局、売上も利益も\n変わらず投資対効果\n不明のまま", COLOR_RED),
]

for i, (title, desc, color) in enumerate(failures):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(3.5), RGBColor(0x24, 0x24, 0x3E), 0.05)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.4), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, x + Inches(0.9), y + Inches(0.55), Inches(1.0), Inches(0.7),
                str(i + 1), font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.4), Inches(0.6),
                title, font_size=22, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(1.0),
                desc, font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "DXの失敗の多くは「ツールの問題」ではなく「運用と測定の問題」",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: DX × 受験アナロジー
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DXを「受験」に例えると", font_size=36, bold=True)

# Left side: 受験
left_card = add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), RGBColor(0x24, 0x24, 0x3E), 0.04)
add_textbox(slide, Inches(1.2), Inches(1.8), Inches(5.0), Inches(0.6),
            "受験の世界", font_size=26, bold=True, color=COLOR_PURPLE)

exam_items = [
    ("合格（ゴール）", "志望校への合格", COLOR_GREEN),
    ("テキスト", "参考書・問題集", COLOR_ACCENT2),
    ("模試・定期テスト", "実力を定期的に測定", COLOR_ORANGE),
    ("学習管理", "日々の勉強習慣の継続", COLOR_PINK),
]

for i, (label, desc, color) in enumerate(exam_items):
    y = Inches(2.6 + i * 0.9)
    badge = add_shape_bg(slide, Inches(1.2), y, Inches(2.4), Inches(0.6), color, 0.3)
    add_textbox(slide, Inches(1.3), y + Inches(0.07), Inches(2.2), Inches(0.45),
                label, font_size=15, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.9), y + Inches(0.07), Inches(2.4), Inches(0.45),
                desc, font_size=15, color=COLOR_LIGHT_GRAY)

# Right side: DX
right_card = add_shape_bg(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0), RGBColor(0x24, 0x24, 0x3E), 0.04)
add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5.0), Inches(0.6),
            "DXの世界", font_size=26, bold=True, color=COLOR_ACCENT)

dx_items = [
    ("合格（ゴール）", "売上・利益の向上", COLOR_GREEN),
    ("DXツール", "SaaS・クラウド・AI", COLOR_ACCENT2),
    ("効果測定", "KPIの定期モニタリング", COLOR_ORANGE),
    ("運用定着", "日々のツール活用と改善", COLOR_PINK),
]

for i, (label, desc, color) in enumerate(dx_items):
    y = Inches(2.6 + i * 0.9)
    badge = add_shape_bg(slide, Inches(7.3), y, Inches(2.4), Inches(0.6), color, 0.3)
    add_textbox(slide, Inches(7.4), y + Inches(0.07), Inches(2.2), Inches(0.45),
                label, font_size=15, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.0), y + Inches(0.07), Inches(2.4), Inches(0.45),
                desc, font_size=15, color=COLOR_LIGHT_GRAY)

# Arrow between
arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(6.1), Inches(3.8), Inches(1.1), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLOR_ACCENT2
arrow.line.fill.background()

# Key insight
key_box = add_shape_bg(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.8), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.1), Inches(9.5), Inches(0.6),
            "テキスト（ツール）を買っただけでは合格（売上向上）できない。模試（効果測定）と学習管理（運用定着）が不可欠。",
            font_size=16, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 4: DX × ダイエットアナロジー
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DXを「ダイエット」に例えると", font_size=36, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "受験は合否（0 or 1）だが、ダイエットは数値の継続改善 ― DXのKPIにより近い",
            font_size=16, color=COLOR_ACCENT2)

# Comparison table style
comparisons = [
    ("ゴール", "体重の削減・体型改善", "売上向上・コスト削減", COLOR_GREEN),
    ("ツール", "ランニングマシン\nトレーニンググッズ", "SaaS・RPA\nクラウドツール", COLOR_ACCENT),
    ("測定", "体重計・体組成計\n定期的な計測", "KPIダッシュボード\n月次レビュー", COLOR_ORANGE),
    ("継続の仕組み", "トレーナー・記録アプリ\n仲間との共有", "推進チーム・定例会議\n利用状況の可視化", COLOR_PURPLE),
    ("よくある失敗", "器具を買って満足\n三日坊主で挫折", "ツールを導入して満足\n現場に定着しない", COLOR_RED),
]

for i, (category, diet, dx, color) in enumerate(comparisons):
    y = Inches(1.9 + i * 1.0)
    # Category label
    cat_badge = add_shape_bg(slide, Inches(0.8), y, Inches(2.0), Inches(0.7), color, 0.3)
    add_textbox(slide, Inches(0.9), y + Inches(0.1), Inches(1.8), Inches(0.5),
                category, font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
    # Diet column
    diet_box = add_shape_bg(slide, Inches(3.1), y, Inches(4.5), Inches(0.7), RGBColor(0x24, 0x24, 0x3E), 0.05)
    add_textbox(slide, Inches(3.3), y + Inches(0.05), Inches(4.1), Inches(0.6),
                diet, font_size=13, color=COLOR_LIGHT_GRAY)
    # DX column
    dx_box = add_shape_bg(slide, Inches(7.9), y, Inches(4.5), Inches(0.7), RGBColor(0x24, 0x24, 0x3E), 0.05)
    add_textbox(slide, Inches(8.1), y + Inches(0.05), Inches(4.1), Inches(0.6),
                dx, font_size=13, color=COLOR_LIGHT_GRAY)

# Column headers
add_textbox(slide, Inches(3.1), Inches(1.55), Inches(4.5), Inches(0.35),
            "ダイエット", font_size=18, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.9), Inches(1.55), Inches(4.5), Inches(0.35),
            "DX", font_size=18, bold=True, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

# Bottom insight
key_box = add_shape_bg(slide, Inches(1.5), Inches(7.2 - 0.3), Inches(10.3), Inches(0.6), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(7.0 - 0.3), Inches(9.5), Inches(0.5),
            "「買って満足」― ダイエット器具もDXツールも、使い続けなければただの置き物",
            font_size=16, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: DX成功の3要素
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DX成功の方程式 ― 3つの要素", font_size=36, bold=True)

elements = [
    ("01", "ツール導入", "Tool Implementation",
     "適切なDXツールの選定と導入",
     ["業務課題に合ったツール選定",
      "既存システムとの連携設計",
      "現場の声を反映した要件定義"],
     COLOR_ACCENT,
     "受験: テキスト購入\nダイエット: 器具購入"),
    ("02", "効果測定", "Measurement & Monitoring",
     "定期的なKPIモニタリングと実力判定",
     ["売上・利益への影響を数値で追跡",
      "ツール利用率・活用度の可視化",
      "月次/四半期での定点観測レビュー"],
     COLOR_ORANGE,
     "受験: 模試・定期テスト\nダイエット: 体重測定"),
    ("03", "運用定着", "Operational Excellence",
     "日々の活用を習慣化し怠けを防ぐ仕組み",
     ["利用状況の自動アラート設定",
      "推進チームによる伴走サポート",
      "成功事例の社内共有と表彰"],
     COLOR_GREEN,
     "受験: 学習計画・管理\nダイエット: トレーナー"),
]

for i, (num, title, en, desc, items, color, analogy) in enumerate(elements):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.6)
    # Card
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(5.3), RGBColor(0x24, 0x24, 0x3E), 0.04)
    # Number
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                num, font_size=14, color=color, bold=True)
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(0.6), Inches(3.1), Inches(0.6),
                title, font_size=28, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.4),
                en, font_size=12, color=COLOR_GRAY)
    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.1), Inches(0.6),
                desc, font_size=14, bold=True, color=COLOR_LIGHT_GRAY)
    # Items
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(2.2 + j * 0.55), Inches(3.1), Inches(0.5),
                    f"• {item}", font_size=13, color=COLOR_LIGHT_GRAY)
    # Analogy box
    analogy_box = add_shape_bg(slide, x + Inches(0.2), y + Inches(4.1), Inches(3.3), Inches(0.9), RGBColor(0x0A, 0x2A, 0x3E), 0.05)
    add_textbox(slide, x + Inches(0.35), y + Inches(4.2), Inches(3.0), Inches(0.7),
                analogy, font_size=11, color=COLOR_GRAY)


# ============================================================
# Slide 6: なぜ「測定」と「運用」が軽視されるのか
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "なぜ「測定」と「運用」が軽視されるのか", font_size=36, bold=True)

reasons = [
    ("導入が「成果」に\n見えてしまう",
     "IR報告や社内報告で「○○を導入しました」\nと言えば、それ自体が実績に見える。\n導入＝ゴールという錯覚が生まれやすい。",
     COLOR_ORANGE),
    ("測定は地味で\n面倒な作業",
     "KPI設計、データ収集、分析、レビュー。\n華やかさがなく、コストもかかる。\n「後でやろう」と後回しにされがち。",
     COLOR_RED),
    ("人は本質的に\n怠ける生き物",
     "受験生もダイエッターも同じ。\n監視と動機づけの仕組みがなければ、\n日々の運用は自然と形骸化する。",
     COLOR_PURPLE),
]

for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(4.0), RGBColor(0x24, 0x24, 0x3E), 0.04)
    # Color top bar
    add_shape_bg(slide, x, y, Inches(3.7), Pt(6), color, 0)
    # Number
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.5), y + Inches(0.3), Inches(0.7), Inches(0.7))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    add_textbox(slide, x + Inches(1.5), y + Inches(0.35), Inches(0.7), Inches(0.55),
                str(i + 1), font_size=24, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.8),
                title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Desc
    add_textbox(slide, x + Inches(0.3), y + Inches(2.2), Inches(3.1), Inches(1.5),
                desc, font_size=14, color=COLOR_LIGHT_GRAY)

# Bottom message
msg_box = add_shape_bg(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.3), Inches(9.5), Inches(0.6),
            "受験でテキストを買って満足する人、ダイエット器具を買って満足する人 ― DXも全く同じ構造",
            font_size=16, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 7: DX成功のためのモニタリング体制
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "DX成功のためのモニタリング体制", font_size=36, bold=True)

# Flow: 導入 → 測定 → 改善 → 定着 (cycle)
monitoring_steps = [
    ("KPI設定", "売上・利益への影響を\n測定可能な指標に落とす", COLOR_ACCENT),
    ("定期測定", "週次/月次で\nダッシュボードを確認", COLOR_GREEN),
    ("利用状況の\n可視化", "誰が・どれだけ\nツールを使っているか", COLOR_ORANGE),
    ("改善アクション", "データに基づく\n運用改善を実行", COLOR_PURPLE),
    ("効果の報告", "経営層へ数値で\n投資対効果を報告", COLOR_ACCENT2),
]

for i, (title, desc, color) in enumerate(monitoring_steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(1.8)
    box = add_shape_bg(slide, x, y, Inches(2.2), Inches(2.5), RGBColor(0x24, 0x24, 0x3E), 0.05)
    add_shape_bg(slide, x, y, Inches(2.2), Pt(6), color, 0)
    # Step number
    add_textbox(slide, x + Inches(0.1), y + Inches(0.2), Inches(0.5), Inches(0.4),
                str(i + 1), font_size=14, color=color, bold=True)
    # Title
    add_textbox(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.0), Inches(0.7),
                title, font_size=18, bold=True, alignment=PP_ALIGN.CENTER, color=color)
    # Desc
    add_textbox(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.0), Inches(1.0),
                desc, font_size=13, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)
    # Arrow
    if i < len(monitoring_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.25), y + Inches(0.95), Inches(0.2), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_ACCENT2
        arrow.line.fill.background()

# Analogy comparison boxes
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.4),
            "アナロジーで理解する「モニタリングの本質」", font_size=20, bold=True, color=COLOR_ACCENT2)

analogies = [
    ("受験", "模試を受けずに本番に臨む受験生はいない。\n定期テストで弱点を把握し、学習計画を修正する。", COLOR_PURPLE),
    ("ダイエット", "体重計に乗らないダイエットは成功しない。\n毎日の計測と記録が行動変容の原動力になる。", COLOR_GREEN),
    ("DX", "KPIを追わないDXは投資の垂れ流し。\n定期的な効果測定が改善サイクルを回す。", COLOR_ACCENT),
]

for i, (label, desc, color) in enumerate(analogies):
    x = Inches(0.8 + i * 4.1)
    y = Inches(5.2)
    a_box = add_shape_bg(slide, x, y, Inches(3.7), Inches(1.6), RGBColor(0x24, 0x24, 0x3E), 0.04)
    add_shape_bg(slide, x, y, Inches(3.7), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.15), Inches(1.5), Inches(0.35),
                label, font_size=18, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.2), Inches(0.9),
                desc, font_size=13, color=COLOR_LIGHT_GRAY)


# ============================================================
# Slide 8: まとめ ― DX成功の方程式
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1.0),
            "まとめ：DX成功の方程式", font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

# Formula visual
formula_box = add_shape_bg(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.8), RGBColor(0x0A, 0x2A, 0x3E), 0.03)

formula_parts = [
    ("DXツール導入", COLOR_ACCENT),
    ("+", COLOR_WHITE),
    ("効果測定", COLOR_ORANGE),
    ("+", COLOR_WHITE),
    ("運用定着", COLOR_GREEN),
    ("=", COLOR_WHITE),
    ("売上・利益向上", COLOR_PINK),
]

for i, (text, color) in enumerate(formula_parts):
    x = Inches(1.3 + i * 1.55)
    if text in ("+", "="):
        add_textbox(slide, x, Inches(2.4), Inches(0.6), Inches(0.8),
                    text, font_size=36, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    else:
        part_box = add_shape_bg(slide, x, Inches(2.3), Inches(1.4), Inches(1.0), RGBColor(0x24, 0x24, 0x3E), 0.1)
        add_textbox(slide, x + Inches(0.05), Inches(2.45), Inches(1.3), Inches(0.7),
                    text, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)

# Key takeaways
takeaways = [
    ("DXツールはあくまで「手段」", "テキストやランニングマシンと同じ。\n持っているだけでは何も変わらない。", COLOR_ACCENT),
    ("定期測定なくして改善なし", "模試や体重計がなければ、\n自分の現在地がわからない。", COLOR_ORANGE),
    ("「怠け」を前提とした仕組みづくり", "人は放っておけば怠ける。\nモニタリングと伴走が必要。", COLOR_GREEN),
    ("最終ゴールは売上と利益", "IR向けの「導入実績」ではなく、\n実際のビジネス成果で評価する。", COLOR_PINK),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.8 + i * 3.1)
    y = Inches(4.3)
    card = add_shape_bg(slide, x, y, Inches(2.8), Inches(2.2), RGBColor(0x24, 0x24, 0x3E), 0.04)
    add_shape_bg(slide, x, y, Inches(2.8), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(1.0),
                desc, font_size=13, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Final message
msg_box = add_shape_bg(slide, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.5), RGBColor(0x0A, 0x2A, 0x3E), 0.03)
add_textbox(slide, Inches(2.8), Inches(6.85), Inches(7.7), Inches(0.4),
            "DXの本当の勝負は、導入の「後」にある。",
            font_size=20, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)


# --- Save ---
output_path = "/home/user/harmonic-marketing/dx_success_metrics.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
