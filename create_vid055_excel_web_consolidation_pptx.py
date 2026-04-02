#!/usr/bin/env python3
"""VID-055: 建設業に必要なのはAIではない - PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xC7)
COLOR_ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
COLOR_CARD = RGBColor(0x24, 0x24, 0x3E)
COLOR_HIGHLIGHT_BG = RGBColor(0x0A, 0x2A, 0x3E)


def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=COLOR_WHITE, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_accent_line(slide, left, top, width, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: タイトル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "建設業に必要なのはAIではない", font_size=48, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
            "まず点在するExcelをWebに統合せよ",
            font_size=24, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.1), Inches(3.3), COLOR_ACCENT)

add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.5),
            "― 建設業DXの現実シリーズ ―",
            font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: AIが使えない本当の理由
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "AIが使えないのは、AIの問題ではない", font_size=36, bold=True)

# IPO flow diagram
ipo_labels = [
    ("Input", "データ入力", COLOR_ACCENT),
    ("Process", "AI処理", COLOR_PURPLE),
    ("Output", "出力・結果", COLOR_GREEN),
]

for i, (label, desc, color) in enumerate(ipo_labels):
    x = Inches(1.5 + i * 3.8)
    y = Inches(1.8)
    box = add_shape_bg(slide, x, y, Inches(2.8), Inches(1.5), COLOR_CARD, 0.05)
    add_shape_bg(slide, x, y, Inches(2.8), Pt(6), color, 0)
    add_textbox(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.5),
                label, font_size=24, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.8), Inches(2.8), Inches(0.5),
                desc, font_size=14, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    if i < 2:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(2.9), y + Inches(0.45),
                                       Inches(0.8), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_ACCENT2
        arrow.line.fill.background()

# Problem highlight
problem_box = add_shape_bg(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.8),
                           RGBColor(0x3E, 0x1A, 0x1A), 0.04)
add_textbox(slide, Inches(1.1), Inches(4.0), Inches(5.0), Inches(0.5),
            "Inputの現状", font_size=20, bold=True, color=COLOR_RED)
problems = [
    "データがバラバラ",
    "フォーマットが統一されていない",
    "最新版がどれかわからない",
    "手入力の転記で整合性がない",
]
for i, p in enumerate(problems):
    add_textbox(slide, Inches(1.3), Inches(4.6 + i * 0.5), Inches(5.0), Inches(0.4),
                f"✕  {p}", font_size=15, color=COLOR_ORANGE)

# Conclusion
conclusion_box = add_shape_bg(slide, Inches(7.0), Inches(3.8), Inches(5.5), Inches(2.8),
                              COLOR_HIGHLIGHT_BG, 0.04)
add_textbox(slide, Inches(7.3), Inches(4.0), Inches(5.0), Inches(0.5),
            "本質", font_size=20, bold=True, color=COLOR_ACCENT)
add_multi_text(slide, Inches(7.3), Inches(4.6), Inches(5.0), Inches(2.0),
               ["「AIが使えない」のではない",
                "",
                "「AIに渡せる状態の",
                "  データがない」のだ"],
               font_size=18, color=COLOR_WHITE, line_spacing=1.6)


# ============================================================
# Slide 3: 業務は手順で動いている
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "業務は「一定の手順」で動いている。しかし……", font_size=36, bold=True)

# Ideal flow
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
            "基本の流れ", font_size=18, bold=True, color=COLOR_GREEN)

flow_steps = ["見積作成", "承認取得", "契約", "着工", "検収", "請求"]
for i, step in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.0)
    y = Inches(2.0)
    box = add_shape_bg(slide, x, y, Inches(1.6), Inches(0.7),
                        RGBColor(0x0A, 0x2E, 0x1A), 0.05)
    add_textbox(slide, x, y + Inches(0.1), Inches(1.6), Inches(0.5),
                step, font_size=14, bold=True, alignment=PP_ALIGN.CENTER,
                color=COLOR_GREEN)
    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(1.62), y + Inches(0.15),
                                       Inches(0.35), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_GREEN
        arrow.line.fill.background()

# Reality
add_textbox(slide, Inches(0.8), Inches(3.2), Inches(3), Inches(0.4),
            "しかし現実は……", font_size=18, bold=True, color=COLOR_RED)

reality_items = [
    ("Aさんの見積", "≠", "Bさんの見積"),
    ("LINE連絡の現場", "≠", "メール連絡の現場"),
    ("部署Aの工程表", "≠", "部署Bの工程表"),
    ("営業の数字", "≠", "経理の数字"),
]

for i, (left, op, right) in enumerate(reality_items):
    y = Inches(3.8 + i * 0.75)
    # Left
    left_box = add_shape_bg(slide, Inches(1.0), y, Inches(3.5), Inches(0.55),
                            COLOR_CARD, 0.05)
    add_textbox(slide, Inches(1.1), y + Inches(0.07), Inches(3.3), Inches(0.4),
                left, font_size=15, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    # ≠
    add_textbox(slide, Inches(4.7), y + Inches(0.05), Inches(0.6), Inches(0.4),
                op, font_size=22, bold=True, color=COLOR_RED,
                alignment=PP_ALIGN.CENTER)
    # Right
    right_box = add_shape_bg(slide, Inches(5.5), y, Inches(3.5), Inches(0.55),
                             COLOR_CARD, 0.05)
    add_textbox(slide, Inches(5.6), y + Inches(0.07), Inches(3.3), Inches(0.4),
                right, font_size=15, color=COLOR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

# Bottom message
msg = add_shape_bg(slide, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.6),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(1.5), Inches(6.65), Inches(10.3), Inches(0.5),
            "手順が整理されていない状態でツールだけ入れても定着しない",
            font_size=16, bold=True, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 4: Excelが点在している現実
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_RED)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "建設業の現実：Excelがあちこちに散らばっている", font_size=36, bold=True)

excel_data = [
    ("見積管理", "見積一覧.xlsx", "営業部の共有フォルダ", COLOR_ACCENT),
    ("工程管理", "工程表.xlsx", "現場所長のPC", COLOR_GREEN),
    ("原価管理", "原価集計.xlsx", "経理のファイルサーバー", COLOR_ORANGE),
    ("日報", "日報_3月.xlsx", "各作業員のスマホ/紙", COLOR_PURPLE),
    ("協力会社管理", "外注先一覧.xlsx", "調達担当の個人フォルダ", COLOR_RED),
    ("発注管理", "発注書.xlsx", "メール添付", COLOR_ACCENT2),
]

for i, (task, filename, location, color) in enumerate(excel_data):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.7)

    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(2.3), COLOR_CARD, 0.04)
    add_shape_bg(slide, x, y, Inches(3.7), Pt(6), color, 0)

    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
                task, font_size=20, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.1), Inches(0.4),
                filename, font_size=14, color=COLOR_LIGHT_GRAY)

    # Location with folder icon
    loc_bg = add_shape_bg(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.1), Inches(0.55),
                          RGBColor(0x1A, 0x1A, 0x2E), 0.05)
    add_textbox(slide, x + Inches(0.4), y + Inches(1.45), Inches(2.9), Inches(0.45),
                location, font_size=13, color=COLOR_GRAY)

# Bottom question
msg = add_shape_bg(slide, Inches(1.5), Inches(7.0 - 0.8), Inches(10.3), Inches(0.6),
                   RGBColor(0x3E, 0x1A, 0x1A), 0.03)
add_textbox(slide, Inches(2.0), Inches(7.0 - 0.75), Inches(9.3), Inches(0.5),
            "この状態でAIを入れて、何ができるでしょうか？",
            font_size=18, bold=True, color=COLOR_RED, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: Excel統合の5ステップ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "最初の一歩：Excel統合の5ステップ", font_size=36, bold=True)

steps = [
    ("01", "業務を1つ選ぶ", "最も痛みが大きい\n業務から着手\n（例：原価管理）"),
    ("02", "Excelを洗い出す", "その業務で使われて\nいるExcelを\n全て列挙する"),
    ("03", "項目を整理する", "項目・入力者\n更新頻度・利用者\nを可視化する"),
    ("04", "重複を特定する", "重複項目・手入力の\n転記箇所を\n見つけ出す"),
    ("05", "Webに集約する", "1つのWebフォーム\n/画面に統合する"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.4 + i * 2.55)
    y = Inches(1.8)

    # Card
    card = add_shape_bg(slide, x, y, Inches(2.3), Inches(4.2), COLOR_CARD, 0.04)
    color = COLOR_GREEN

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x + Inches(0.8), y + Inches(0.3),
                                   Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_textbox(slide, x + Inches(0.8), y + Inches(0.35), Inches(0.7), Inches(0.6),
                num, font_size=22, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(1.9), Inches(0.6),
                title, font_size=18, bold=True, color=COLOR_ACCENT2,
                alignment=PP_ALIGN.CENTER)

    # Description
    add_multi_text(slide, x + Inches(0.2), y + Inches(1.9), Inches(1.9), Inches(2.0),
                   desc.split("\n"), font_size=13, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.5)

    # Arrow
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(2.35), y + Inches(1.8),
                                       Inches(0.18), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_GREEN
        arrow.line.fill.background()

# Result message
result = add_shape_bg(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8),
                      COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(6.4), Inches(9.3), Inches(0.6),
            "統合されたデータがあれば、集計も分析もAI活用も、初めて現実的になる",
            font_size=17, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 6: 見るべき指標
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "見るべき指標は決まっている。なのに共有されていない。", font_size=34, bold=True)

metrics = [
    ("完成工事粗利率", "案件ごとに\n利益が出ているか", COLOR_GREEN),
    ("外注費比率", "自社施工と外注の\nバランスは適正か", COLOR_ACCENT),
    ("進捗率 vs 原価消化率", "コスト超過の\n兆候はないか", COLOR_ORANGE),
    ("変更契約の回収率", "追加工事の\n請求漏れはないか", COLOR_RED),
    ("未成工事受入金", "資金繰りは\n大丈夫か", COLOR_PURPLE),
]

for i, (title, desc, color) in enumerate(metrics):
    x = Inches(0.4 + i * 2.55)
    y = Inches(1.6)

    card = add_shape_bg(slide, x, y, Inches(2.3), Inches(2.5), COLOR_CARD, 0.04)
    add_shape_bg(slide, x, y, Inches(2.3), Pt(6), color, 0)

    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(2.0), Inches(0.7),
                title, font_size=16, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.0), Inches(1.0),
                   desc.split("\n"), font_size=13, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.4)

# Problem section
prob_box = add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
                        RGBColor(0x3E, 0x1A, 0x1A), 0.04)
add_textbox(slide, Inches(1.1), Inches(4.7), Inches(5.0), Inches(0.4),
            "現状の問題", font_size=18, bold=True, color=COLOR_RED)
add_multi_text(slide, Inches(1.1), Inches(5.2), Inches(5.0), Inches(1.5),
               ["「あの部長が見ている」",
                "「経理のCさんしかわからない」",
                "→ 属人的な管理。ノウハウが共有されない。"],
               font_size=15, color=COLOR_ORANGE, line_spacing=1.6)

# Consultant section
cons_box = add_shape_bg(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5),
                        COLOR_HIGHLIGHT_BG, 0.04)
add_textbox(slide, Inches(7.3), Inches(4.7), Inches(5.0), Inches(0.4),
            "コンサルタントの問題", font_size=18, bold=True, color=COLOR_ACCENT)
add_multi_text(slide, Inches(7.3), Inches(5.2), Inches(5.0), Inches(1.5),
               ["「お客様が見たいもの」を見せるだけ",
                "見るべき指標を定義できていない",
                "→ 御用聞きはコンサルティングではない"],
               font_size=15, color=COLOR_LIGHT_GRAY, line_spacing=1.6)


# ============================================================
# Slide 7: 正しいDX順序
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
            "正しい順序で進める", font_size=36, bold=True)

dx_steps = [
    ("①", "業務整理", "手順を可視化\n標準化する", COLOR_ACCENT),
    ("②", "情報統合", "Excel→Webに\n集約する", COLOR_ACCENT2),
    ("③", "指標定義", "見るべき数字を\n決める", COLOR_GREEN),
    ("④", "ツール選定", "業務に合うツール\nを選ぶ", COLOR_ORANGE),
    ("⑤", "AI活用", "統合データで\nAIを活用", COLOR_PURPLE),
]

# Timeline bar
add_shape_bg(slide, Inches(1.2), Inches(2.5), Inches(10.9), Inches(0.15),
             COLOR_ACCENT, 0.5)

for i, (num, title, desc, color) in enumerate(dx_steps):
    x = Inches(0.7 + i * 2.45)

    # Circle on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + Inches(0.8), Inches(2.25),
                                 Inches(0.55), Inches(0.55))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_textbox(slide, x + Inches(0.8), Inches(2.3), Inches(0.55), Inches(0.45),
                num, font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

    # Card below
    card = add_shape_bg(slide, x, Inches(3.2), Inches(2.2), Inches(2.5),
                        COLOR_CARD, 0.04)
    add_shape_bg(slide, x, Inches(3.2), Inches(2.2), Pt(6), color, 0)

    add_textbox(slide, x + Inches(0.1), Inches(3.5), Inches(2.0), Inches(0.5),
                title, font_size=20, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.1), Inches(4.1), Inches(2.0), Inches(1.2),
                   desc.split("\n"), font_size=13, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.4)

# Warning box
warn = add_shape_bg(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.0),
                    RGBColor(0x3E, 0x1A, 0x1A), 0.03)
add_textbox(slide, Inches(2.0), Inches(6.1), Inches(3.0), Inches(0.4),
            "⚠ 多くの会社が犯す間違い", font_size=16, bold=True, color=COLOR_RED)
add_textbox(slide, Inches(2.0), Inches(6.5), Inches(9.5), Inches(0.4),
            "①②③を飛ばして④⑤から始める → AIは万能ではない。整理されたデータの上でだけ力を発揮する。",
            font_size=15, bold=True, color=COLOR_ORANGE)


# ============================================================
# Slide 8: まとめ・CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), prs.slide_width, COLOR_ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1.0),
            "まとめ：まずExcelを集めることから",
            font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

# 3 takeaways
takeaways = [
    ("AIの前にデータ統合", "AIに渡せる状態の\nデータを作ることが先", COLOR_ACCENT),
    ("Excel統合5ステップ", "痛みが大きい業務から\n1つずつ集約していく", COLOR_GREEN),
    ("見るべき指標を定義", "属人管理から脱却し\n標準指標を全社共有", COLOR_PURPLE),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.3)
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(2.2), COLOR_CARD, 0.04)
    add_shape_bg(slide, x, y, Inches(3.7), Pt(5), color, 0)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.4), Inches(3.1), Inches(0.6),
                title, font_size=22, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.8),
                   desc.split("\n"), font_size=15, color=COLOR_LIGHT_GRAY,
                   line_spacing=1.4)

# Key message
key = add_shape_bg(slide, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.8),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.8), Inches(4.9), Inches(7.7), Inches(0.6),
            "地味な一歩。しかしこの一歩なしにDXは進まない。",
            font_size=20, bold=True, color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

# CTA box
cta = add_shape_bg(slide, Inches(1.5), Inches(5.9), Inches(10.3), Inches(1.3),
                   COLOR_HIGHLIGHT_BG, 0.03)
add_textbox(slide, Inches(2.0), Inches(6.0), Inches(9), Inches(0.4),
            "HARMONIC insight ── 建設業のExcel統合からDX実現まで",
            font_size=18, bold=True, color=COLOR_ACCENT)

cta_items = [
    "自分で始めたい方 → note記事「建設業に必要なのはAIではない」",
    "体系的に学びたい方 → SIPOフレームワーク実践ワークショップ",
    "支援が必要な方 → お問い合わせ（概要欄リンク）",
]
for i, item in enumerate(cta_items):
    add_textbox(slide, Inches(2.5), Inches(6.4 + i * 0.35), Inches(8.5), Inches(0.35),
                f"▸ {item}", font_size=13, color=COLOR_LIGHT_GRAY)


# --- Save ---
output_path = "/home/user/harmonic-marketing/content/youtube/slides/VID-055_excel_web_consolidation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
