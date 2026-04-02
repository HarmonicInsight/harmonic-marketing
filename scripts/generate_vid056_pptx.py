#!/usr/bin/env python3
"""Generate VID-056 YouTube presentation: AI会計自動化の落とし穴"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Widescreen 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)


def add_dark_bg(slide):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text_box(slide, left, top, width, height, text, font_size=24,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=22,
                       color=WHITE, font_name="Meiryo", line_spacing=1.5):
    """Add multi-line text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bld if bld else False
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=18, font_color=WHITE):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Meiryo"
    return shape


# ─── Slide 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.5, 12, 1, "VID-056 / HARMONIC insight", 16, MID_GRAY)
add_text_box(slide, 0.8, 2.0, 11.5, 2.0,
             "AIで数字が出る ＝ 正しい、ではない", 48, WHITE, bold=True)
add_text_box(slide, 0.8, 4.0, 11.5, 1.5,
             "AI会計自動化の落とし穴\n── 監査を通すために、先にやるべきこと", 28, ACCENT_BLUE)
add_text_box(slide, 0.8, 6.5, 6, 0.5, "HARMONIC insight / 瀬田裕之", 18, LIGHT_GRAY)

# ─── Slide 2: Problem Statement ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "何が起きているのか", 36, ACCENT_BLUE, bold=True)

add_multiline_text(slide, 0.8, 1.3, 11.5, 2.5, [
    ("Gemini × Googleスプレッドシートで、自然言語から数式を自動生成", 24, WHITE, False),
    ("", 12, WHITE, False),
    ("「A列とB列から合計を出して」→ AIが =A1+B1 を自動入力", 22, LIGHT_GRAY, False),
    ("「関数知識が不要になる」「生産性格差が縮小する」", 22, LIGHT_GRAY, False),
])

add_rounded_rect(slide, 1.5, 4.5, 10, 2.0, RGBColor(0x33, 0x22, 0x22),
                 "", 24, WHITE)
add_text_box(slide, 2.0, 4.7, 9, 0.8,
             "でも、会計の世界では", 24, LIGHT_GRAY)
add_text_box(slide, 2.0, 5.3, 9, 1.0,
             "「数字が出る」と「数字が正しい」はまったく別の話", 30, ACCENT_RED, bold=True)

# ─── Slide 3: What audits require ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "監査で問われる5つの要件", 36, ACCENT_BLUE, bold=True)

items = [
    ("算出根拠", "なぜその計算式なのか。どの基準に基づくのか"),
    ("計算手順", "どのデータを、どの順番で、どう処理したか"),
    ("承認フロー", "誰が確認し、誰が承認したか"),
    ("変更履歴", "いつ、誰が、何を変更したか"),
    ("監査証跡", "上記すべてが追跡可能であること"),
]
for i, (title, desc) in enumerate(items):
    y = 1.5 + i * 1.1
    add_rounded_rect(slide, 0.8, y, 2.5, 0.8, ACCENT_BLUE, title, 20, WHITE)
    add_text_box(slide, 3.6, y + 0.1, 9, 0.7, desc, 20, LIGHT_GRAY)

add_text_box(slide, 0.8, 7.0, 12, 0.5,
             "AIが数式を自動生成すると、これらの要件が満たされない", 20, ACCENT_RED, bold=True)

# ─── Slide 4: The dangerous equation ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "危険な等式", 36, ACCENT_BLUE, bold=True)

pairs = [
    ("数字が出る", "≠", "数字が正しい"),
    ("計算できる", "≠", "根拠がある"),
    ("自動化された", "≠", "統制されている"),
]
for i, (left, eq, right) in enumerate(pairs):
    y = 2.0 + i * 1.5
    add_rounded_rect(slide, 1.5, y, 3.5, 0.9, RGBColor(0x33, 0x33, 0x55), left, 26, WHITE)
    add_text_box(slide, 5.2, y, 1.5, 0.9, eq, 40, ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, 7.0, y, 4.5, 0.9, RGBColor(0x33, 0x33, 0x55), right, 26, WHITE)

add_text_box(slide, 1.5, 6.5, 10, 0.8,
             "AIは計算を速くする道具であって、算出根拠を作る道具ではない",
             24, ACCENT_YELLOW, bold=True)

# ─── Slide 5: Construction industry risks ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "建設業で特にリスクが高い3つの理由", 36, ACCENT_BLUE, bold=True)

risks = [
    ("1", "工事進行基準", "見積総原価の妥当性は人間が判断する必要がある\n当期売上 = 契約金額 x 工事進捗度"),
    ("2", "工事別原価管理", "配賦基準を定義しないまま「原価を計算して」では\n正しい数字は出ない"),
    ("3", "完成工事原価報告書", "合計だけでなく内訳の積み上げ根拠が必要\n税務申告に耐える精度が求められる"),
]
for i, (num, title, desc) in enumerate(risks):
    y = 1.5 + i * 1.8
    add_rounded_rect(slide, 0.8, y, 0.8, 0.8, ACCENT_RED, num, 28, WHITE)
    add_text_box(slide, 1.8, y, 3, 0.8, title, 26, WHITE, bold=True)
    add_text_box(slide, 1.8, y + 0.7, 10, 1.0, desc, 18, LIGHT_GRAY)

# ─── Slide 6: Wrong vs Right order ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "順序を間違えると、コストが何倍にもなる", 34, ACCENT_BLUE, bold=True)

# Wrong order
add_text_box(slide, 0.8, 1.3, 6, 0.6, "やってはいけない順序", 22, ACCENT_RED, bold=True)
wrong_steps = ["AIツール導入", "数字が出る", "「便利だ」", "監査で指摘", "根拠の後付け", "大幅な手戻り"]
for i, step in enumerate(wrong_steps):
    x = 0.8 + i * 2.0
    clr = ACCENT_RED if i >= 3 else RGBColor(0x55, 0x44, 0x44)
    add_rounded_rect(slide, x, 2.0, 1.8, 0.7, clr, step, 14, WHITE)
    if i < len(wrong_steps) - 1:
        add_text_box(slide, x + 1.8, 2.0, 0.3, 0.7, "→", 18, MID_GRAY, alignment=PP_ALIGN.CENTER)

# Right order
add_text_box(slide, 0.8, 3.5, 6, 0.6, "正しい順序", 22, ACCENT_GREEN, bold=True)
right_steps = [
    ("Step 1", "算出基準の定義"),
    ("Step 2", "計算手順の文書化"),
    ("Step 3", "承認フロー・\n監査証跡の設計"),
    ("Step 4", "AI自動化を適用"),
]
for i, (num, step) in enumerate(right_steps):
    x = 0.8 + i * 3.0
    add_rounded_rect(slide, x, 4.2, 2.7, 1.2, ACCENT_GREEN if i < 3 else ACCENT_BLUE,
                     f"{num}\n{step}", 16, WHITE)
    if i < len(right_steps) - 1:
        add_text_box(slide, x + 2.7, 4.5, 0.4, 0.7, "→", 22, MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.0, 12, 1.0,
             "Step 1〜3を飛ばしてAIを入れると、後からの整備コストが導入コストの何倍にもなる",
             22, ACCENT_YELLOW, bold=True)

# ─── Slide 7: What to do concretely ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "具体的に何をすべきか", 36, ACCENT_BLUE, bold=True)

actions = [
    ("算出基準書の作成", [
        "計算式の明文化（例：粗利 = 契約金額 - 実行予算原価合計）",
        "準拠基準の記載（会計基準・社内規程）",
        "入力データの定義・例外処理",
    ]),
    ("計算手順書の整備", [
        "データ取得元と取得方法",
        "検証方法（ダブルチェック・クロスチェック）",
        "承認者と承認基準",
    ]),
    ("AI活用ルールの策定", [
        "AI生成の数式は必ず人間が検証",
        "AI計算結果は「下書き」→ 承認後「確定」",
        "使用した計算ロジックをログに記録",
    ]),
]

for i, (title, items) in enumerate(actions):
    x = 0.8 + i * 4.0
    add_rounded_rect(slide, x, 1.4, 3.7, 0.8, ACCENT_BLUE, title, 18, WHITE)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.1, 2.4 + j * 0.7, 3.6, 0.6,
                     f"  {item}", 14, LIGHT_GRAY)

# ─── Slide 8: Summary / CTA ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, 0.8, 0.3, 12, 1, "まとめ", 36, ACCENT_BLUE, bold=True)

add_multiline_text(slide, 0.8, 1.5, 11.5, 4.0, [
    ("AI会計自動化は「便利」だが、それだけでは監査に耐えられない", 26, WHITE, True),
    ("", 12, WHITE, False),
    ("算出根拠・計算手順・承認フロー・変更履歴 ── すべてが揃って初めて、数字に意味が出る", 22, LIGHT_GRAY, False),
    ("", 12, WHITE, False),
    ("正しい順序：基準定義 → 手順文書化 → 監査証跡設計 → AI自動化", 24, ACCENT_GREEN, True),
    ("", 12, WHITE, False),
    ("この順序を守れる会社だけが、AI自動化の恩恵を本当に受けられる", 24, ACCENT_YELLOW, True),
])

add_rounded_rect(slide, 2.0, 5.5, 9, 1.2, RGBColor(0x22, 0x33, 0x44),
                 "HARMONIC insight ── 管理会計整備からAI活用ルール策定まで一貫支援", 20, ACCENT_BLUE)

add_text_box(slide, 0.8, 7.0, 12, 0.5,
             "チャンネル登録 & 高評価お願いします", 18, MID_GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_dir = "/home/user/harmonic-marketing/content/youtube/slides"
output_path = os.path.join(output_dir, "VID-056_ai-accounting-pitfall.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
