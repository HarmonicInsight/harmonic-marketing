"""
ニュース速報分析シリーズ - Shorts用PPTX生成スクリプト
既存デザイン（ダーク背景 + ゴールドアクセント）を踏襲しつつ、
日付ラベルやデータビジュアルを追加したグラフィカルなスライドを生成。
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 定数 ──
SLIDE_W = Emu(4114800)  # 4.5 inches
SLIDE_H = Emu(7315200)  # 8.0 inches

# カラーパレット
BG_DARK = RGBColor(0x1C, 0x19, 0x17)
GOLD = RGBColor(0xB8, 0x94, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM_GRAY = RGBColor(0x57, 0x53, 0x4E)
LIGHT_GRAY = RGBColor(0xA8, 0xA2, 0x9E)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)
BLUE_ACCENT = RGBColor(0x42, 0xA5, 0xF5)
GREEN_ACCENT = RGBColor(0x66, 0xBB, 0x6A)
ORANGE_ACCENT = RGBColor(0xFF, 0xA7, 0x26)
TEAL_ACCENT = RGBColor(0x26, 0xC6, 0xDA)
DARK_SURFACE = RGBColor(0x2C, 0x27, 0x23)
PURPLE_ACCENT = RGBColor(0xAB, 0x47, 0xBC)

FONT_NAME = "Hiragino Sans"
FONT_NAME_FALLBACK = "Yu Gothic"

MARGIN_X = Emu(274320)
CONTENT_W = Emu(3566160)
TOP_BAR_H = Emu(54864)
BRAND_Y = Emu(6675120)
BRAND_H = Emu(365760)

OUTPUT_DIR = "content/youtube/shorts"


def _set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_top_bar(slide, color=GOLD):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TOP_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_brand(slide):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, BRAND_Y, CONTENT_W, BRAND_H
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HARMONIC insight"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(9)
    run.font.color.rgb = GOLD


def _add_date_label(slide, date_text, source_text, x=None, y=None):
    """赤い日付ラベル + ソース名"""
    lbl_x = x or Emu(137160)
    lbl_y = y or Emu(274320)
    lbl_w = Emu(1920240)
    lbl_h = Emu(365760)

    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, lbl_x, lbl_y, lbl_w, lbl_h
    )
    label.fill.solid()
    label.fill.fore_color.rgb = RED_ACCENT
    label.line.fill.background()
    # Adjust corner rounding
    label.adjustments[0] = 0.25

    tf = label.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = date_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Source name next to label
    src_x = lbl_x + lbl_w + Emu(91440)
    src = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, src_x, lbl_y, Emu(1554480), lbl_h
    )
    src.fill.background()
    src.line.fill.background()
    tf2 = src.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.text = source_text
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = FONT_NAME
    run2.font.size = Pt(11)
    run2.font.color.rgb = LIGHT_GRAY


def _add_text_block(slide, texts, y, h, font_size, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    """メインテキストブロック"""
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y), CONTENT_W, Emu(h)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_sub_text(slide, texts, y, h, font_size=12, color=WARM_GRAY):
    _add_text_block(slide, texts, y, h, font_size, color, bold=False)


def _add_horizontal_bar(slide, label, value, max_val, y, color, width_ratio=0.85):
    """水平バーチャート（1本）"""
    bar_x = MARGIN_X
    bar_y = Emu(y)
    bar_h = Emu(320040)
    max_w = int(CONTENT_W * width_ratio)
    bar_w = int(max_w * (value / max_val))

    # Label
    lbl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_x, Emu(y - 228600), CONTENT_W, Emu(228600)
    )
    lbl.fill.background()
    lbl.line.fill.background()
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(11)
    run.font.color.rgb = LIGHT_GRAY

    # Bar background
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, Emu(max_w), bar_h
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = DARK_SURFACE
    bg_bar.line.fill.background()
    bg_bar.adjustments[0] = 0.15

    # Bar fill
    if bar_w > 0:
        fill_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, Emu(bar_w), bar_h
        )
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = color
        fill_bar.line.fill.background()
        fill_bar.adjustments[0] = 0.15

    # Value
    val_x = Emu(bar_x + max_w + 68580)
    val = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, val_x, bar_y, Emu(548640), bar_h
    )
    val.fill.background()
    val.line.fill.background()
    tf2 = val.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = f"{value}%"
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = FONT_NAME
    run2.font.size = Pt(13)
    run2.font.bold = True
    run2.font.color.rgb = color


def _add_big_number(slide, number, unit, description, y, color=GOLD):
    """大きな数字表示"""
    # Number
    num_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y), CONTENT_W, Emu(640080)
    )
    num_box.fill.background()
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run_num = p.add_run()
    run_num.text = str(number)
    run_num.font.name = FONT_NAME
    run_num.font.size = Pt(52)
    run_num.font.bold = True
    run_num.font.color.rgb = color

    run_unit = p.add_run()
    run_unit.text = unit
    run_unit.font.name = FONT_NAME
    run_unit.font.size = Pt(20)
    run_unit.font.bold = True
    run_unit.font.color.rgb = color

    # Description
    if description:
        desc_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y + 640080), CONTENT_W, Emu(320040)
        )
        desc_box.fill.background()
        desc_box.line.fill.background()
        tf2 = desc_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = description
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0]
        run2.font.name = FONT_NAME
        run2.font.size = Pt(12)
        run2.font.color.rgb = LIGHT_GRAY


def _add_vs_block(slide, left_val, left_label, right_val, right_label, y, left_color, right_color):
    """VS比較ブロック"""
    block_y = Emu(y)
    half_w = int(CONTENT_W / 2) - Emu(45720)

    # Left box
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, block_y, Emu(half_w), Emu(960120)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = DARK_SURFACE
    left_box.line.fill.background()
    left_box.adjustments[0] = 0.08

    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_val = tf_l.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    run_v = p_val.add_run()
    run_v.text = left_val
    run_v.font.name = FONT_NAME
    run_v.font.size = Pt(36)
    run_v.font.bold = True
    run_v.font.color.rgb = left_color

    p_lbl = tf_l.add_paragraph()
    p_lbl.alignment = PP_ALIGN.CENTER
    run_l = p_lbl.add_run()
    run_l.text = left_label
    run_l.font.name = FONT_NAME
    run_l.font.size = Pt(10)
    run_l.font.color.rgb = LIGHT_GRAY

    # VS circle
    vs_x = Emu(int(MARGIN_X + half_w + 45720 / 2 - 182880))
    vs_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, vs_x, Emu(y + 320040), Emu(365760), Emu(365760)
    )
    vs_circle.fill.solid()
    vs_circle.fill.fore_color.rgb = RED_ACCENT
    vs_circle.line.fill.background()
    tf_vs = vs_circle.text_frame
    tf_vs.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_vs = tf_vs.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    run_vs = p_vs.add_run()
    run_vs.text = "VS"
    run_vs.font.name = FONT_NAME
    run_vs.font.size = Pt(11)
    run_vs.font.bold = True
    run_vs.font.color.rgb = WHITE

    # Right box
    right_x = Emu(int(MARGIN_X + half_w + 91440))
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, block_y, Emu(half_w), Emu(960120)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = DARK_SURFACE
    right_box.line.fill.background()
    right_box.adjustments[0] = 0.08

    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_val2 = tf_r.paragraphs[0]
    p_val2.alignment = PP_ALIGN.CENTER
    run_v2 = p_val2.add_run()
    run_v2.text = right_val
    run_v2.font.name = FONT_NAME
    run_v2.font.size = Pt(36)
    run_v2.font.bold = True
    run_v2.font.color.rgb = right_color

    p_lbl2 = tf_r.add_paragraph()
    p_lbl2.alignment = PP_ALIGN.CENTER
    run_l2 = p_lbl2.add_run()
    run_l2.text = right_label
    run_l2.font.name = FONT_NAME
    run_l2.font.size = Pt(10)
    run_l2.font.color.rgb = LIGHT_GRAY


def _add_divider(slide, y, color=GOLD):
    """区切り線"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y), CONTENT_W, Emu(18288)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide)
    _add_top_bar(slide)
    _add_brand(slide)
    return slide


def _add_icon_text_row(slide, icon_char, text, y, icon_color=GOLD, text_color=WHITE):
    """アイコン + テキスト行"""
    # Icon
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, MARGIN_X, Emu(y), Emu(320040), Emu(320040)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = DARK_SURFACE
    icon_box.line.fill.background()
    tf_i = icon_box.text_frame
    tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_i = tf_i.paragraphs[0]
    p_i.alignment = PP_ALIGN.CENTER
    run_i = p_i.add_run()
    run_i.text = icon_char
    run_i.font.name = FONT_NAME
    run_i.font.size = Pt(14)
    run_i.font.bold = True
    run_i.font.color.rgb = icon_color

    # Text
    txt_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(MARGIN_X + 388620), Emu(y), Emu(CONTENT_W - 388620), Emu(320040)
    )
    txt_box.fill.background()
    txt_box.line.fill.background()
    tf_t = txt_box.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]
    p_t.text = text
    run_t = p_t.runs[0]
    run_t.font.name = FONT_NAME
    run_t.font.size = Pt(13)
    run_t.font.color.rgb = text_color


# ============================================================
# PPTX 1: Anthropic 81,000人調査 - 生産性じゃなく自由
# ============================================================
def create_anthropic_freedom():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: タイトル + 日付ラベル ──
    s = _new_slide(prs)
    _add_date_label(s, "3月18日 発表", "Anthropic")
    _add_big_number(s, "81,000", "人", "159か国 · 70言語", 1143000)
    _add_text_block(s, [
        "AIに何を求めるか？",
        "史上最大の調査",
    ], 2514600, 1371600, 28, WHITE, True)
    _add_sub_text(s, [
        "答えは「生産性」じゃなかった",
    ], 4114800, 640080, 15, GOLD)

    # ── Slide 2: 9つのビジョン トップ4 バーチャート ──
    s2 = _new_slide(prs)
    _add_date_label(s2, "3月18日 発表", "Anthropic")
    _add_text_block(s2, [
        "AIに求めるもの",
        "トップ4",
    ], 822960, 731520, 24, WHITE, True)

    _add_horizontal_bar(s2, "職業的卓越性", 18.8, 25, 2057400, BLUE_ACCENT)
    _add_horizontal_bar(s2, "個人的変容", 13.7, 25, 2743200, ORANGE_ACCENT)
    _add_horizontal_bar(s2, "生活管理", 13.5, 25, 3429000, GREEN_ACCENT)
    _add_horizontal_bar(s2, "時間的自由", 11.1, 25, 4114800, TEAL_ACCENT)

    _add_sub_text(s2, [
        "上位3つで46%。どれも単純な",
        "「生産性向上」には収まらない",
    ], 5029200, 640080, 12, WARM_GRAY)

    # ── Slide 3: 深く聞くと変わる ──
    s3 = _new_slide(prs)
    _add_text_block(s3, [
        "「なぜそれが",
        "  欲しいの？」",
        "と深く聞くと…",
    ], 1143000, 1828800, 24, WHITE, True)
    _add_divider(s3, 3200400, GOLD)
    _add_icon_text_row(s3, "14%", "自分自身の成長", 3429000, ORANGE_ACCENT, WHITE)
    _add_icon_text_row(s3, "14%", "認知的負担を減らしたい", 3886200, GREEN_ACCENT, WHITE)
    _add_icon_text_row(s3, "11%", "家族と過ごす時間", 4343400, TEAL_ACCENT, WHITE)

    _add_sub_text(s3, [
        "メキシコのエンジニア：",
        "「AIのおかげで定時に帰れる。",
        "  子どもの迎えに行ける」",
    ], 5029200, 914400, 12, WARM_GRAY)

    # ── Slide 4: 結論 ──
    s4 = _new_slide(prs)
    _add_text_block(s4, [
        "人は",
        "「生産性」と",
        "答えた",
    ], 1371600, 1600200, 28, WHITE, True)
    _add_divider(s4, 3200400, GOLD)
    _add_text_block(s4, [
        "本当に",
        "求めていたのは",
        "「自由」だった",
    ], 3429000, 1600200, 28, GOLD, True)

    # ── Slide 5: 締め ──
    s5 = _new_slide(prs)
    _add_text_block(s5, [
        "生産性は手段。",
        "目的は、",
        "人生を取り戻す",
        "こと。",
    ], 2057400, 2286000, 28, WHITE, True)
    _add_sub_text(s5, [
        "出典: Anthropic \"81,000 Interviews\" 2026",
    ], 5486400, 457200, 9, WARM_GRAY)

    path = os.path.join(OUTPUT_DIR, "SHORT_NEWS_01_anthropic_freedom.pptx")
    prs.save(path)
    print(f"Created: {path}")


# ============================================================
# PPTX 2: Anthropic 自律性ギャップ 47% vs 14%
# ============================================================
def create_anthropic_autonomy():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: タイトル ──
    s = _new_slide(prs)
    _add_date_label(s, "3月18日 発表", "Anthropic")
    _add_text_block(s, [
        "AIの恩恵、",
        "会社員と",
        "フリーランスで",
    ], 1143000, 1600200, 24, WHITE, True)
    _add_text_block(s, [
        "3倍違う",
    ], 2971800, 731520, 36, RED_ACCENT, True)
    _add_sub_text(s, [
        "81,000人調査が示した決定的な差",
    ], 4114800, 457200, 12, WARM_GRAY)

    # ── Slide 2: VS比較 ──
    s2 = _new_slide(prs)
    _add_date_label(s2, "3月18日 発表", "Anthropic")
    _add_text_block(s2, [
        "AIで経済的恩恵を",
        "感じている人の割合",
    ], 822960, 731520, 18, LIGHT_GRAY, False)

    _add_vs_block(s2,
        "47%", "独立事業者", "14%", "企業従業員",
        1828800, GREEN_ACCENT, RED_ACCENT
    )

    _add_text_block(s2, [
        "3.4倍の差",
    ], 3200400, 548640, 24, GOLD, True, PP_ALIGN.CENTER)

    _add_divider(s2, 3886200, GOLD)
    _add_text_block(s2, [
        "副業持ち会社員",
    ], 4114800, 365760, 14, LIGHT_GRAY, False)
    _add_text_block(s2, [
        "58%",
    ], 4480560, 548640, 36, ORANGE_ACCENT, True, PP_ALIGN.CENTER)

    # ── Slide 3: 問題提起 ──
    s3 = _new_slide(prs)
    _add_text_block(s3, [
        "AIの恩恵は",
        "能力で決まらない",
    ], 1371600, 1371600, 24, WHITE, True)
    _add_divider(s3, 2971800, GOLD)
    _add_text_block(s3, [
        "自律性で決まる",
    ], 3200400, 914400, 28, GOLD, True)
    _add_sub_text(s3, [
        "承認プロセス・セキュリティポリシー…",
        "「まず上に確認します」の間に、",
        "フリーランスはAIで終わらせている",
    ], 4571040, 914400, 12, WARM_GRAY)

    # ── Slide 4: 日本企業への示唆 ──
    s4 = _new_slide(prs)
    _add_text_block(s4, [
        "日本企業の",
        "79.3%が",
    ], 1143000, 1143000, 24, WHITE, True)
    _add_text_block(s4, [
        "「人材不足」",
    ], 2286000, 731520, 28, RED_ACCENT, True, PP_ALIGN.CENTER)
    _add_text_block(s4, [
        "と答えた",
    ], 2971800, 548640, 20, WHITE, True, PP_ALIGN.CENTER)
    _add_divider(s4, 3657600, GOLD)
    _add_text_block(s4, [
        "本当の問題は",
        "人材がいないことじゃなく",
        "裁量を渡していないこと",
    ], 3886200, 1371600, 18, GOLD, True)

    # ── Slide 5: 締め ──
    s5 = _new_slide(prs)
    _add_text_block(s5, [
        "自律性を",
        "渡さない組織は、",
        "AIの恩恵を",
        "受けられない。",
    ], 2057400, 2286000, 28, WHITE, True)
    _add_sub_text(s5, [
        "出典: Anthropic \"81,000 Interviews\" 2026",
    ], 5486400, 457200, 9, WARM_GRAY)

    path = os.path.join(OUTPUT_DIR, "SHORT_NEWS_02_anthropic_autonomy.pptx")
    prs.save(path)
    print(f"Created: {path}")


# ============================================================
# PPTX 3: Anthropic 光と影
# ============================================================
def create_anthropic_light_shade():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: タイトル ──
    s = _new_slide(prs)
    _add_date_label(s, "3月18日 発表", "Anthropic")
    _add_text_block(s, [
        "AIを一番",
        "使っている人が、",
    ], 1371600, 1143000, 24, WHITE, True)
    _add_text_block(s, [
        "一番AIを",
        "恐れている",
    ], 2743200, 1143000, 28, RED_ACCENT, True)
    _add_sub_text(s, [
        "81,000人が語った「光と影」",
    ], 4343400, 457200, 13, WARM_GRAY)

    # ── Slide 2: 3倍データ ──
    s2 = _new_slide(prs)
    _add_date_label(s2, "3月18日 発表", "Anthropic")
    _add_text_block(s2, [
        "AIに感情的に",
        "助けられた人は",
    ], 1143000, 914400, 20, WHITE, True)
    _add_text_block(s2, [
        "依存を恐れる率",
    ], 2286000, 548640, 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    _add_big_number(s2, "3", "倍", "", 2743200, RED_ACCENT)

    _add_divider(s2, 3886200, GOLD)

    _add_text_block(s2, [
        "教師が学生の思考力低下を",
        "目撃する率",
    ], 4114800, 731520, 14, LIGHT_GRAY, False)
    _add_text_block(s2, [
        "平均の2.5倍",
    ], 4800600, 548640, 24, ORANGE_ACCENT, True)

    # ── Slide 3: 弁護士の声 ──
    s3 = _new_slide(prs)

    # Quote decoration
    quote_bar = slide = s3
    qbar = s3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(137160), Emu(1371600), Emu(36576), Emu(2057400)
    )
    qbar.fill.solid()
    qbar.fill.fore_color.rgb = GOLD
    qbar.line.fill.background()

    _add_text_block(s3, [
        "AIで契約書レビューの",
        "時間を節約している。",
    ], 1371600, 1143000, 18, WHITE, True)
    _add_text_block(s3, [
        "でも同時に恐れている。",
        "自分で読む力を",
        "失っているのでは",
        "ないか。",
    ], 2514600, 1371600, 18, ORANGE_ACCENT, True)
    _add_sub_text(s3, [
        "── イスラエルの弁護士",
    ], 4114800, 365760, 12, WARM_GRAY)

    # ── Slide 4: 懸念トップ3 ──
    s4 = _new_slide(prs)
    _add_text_block(s4, [
        "懸念トップ3",
    ], 822960, 548640, 22, GOLD, True)

    _add_horizontal_bar(s4, "1. 信頼性（ハルシネーション）", 26.7, 35, 1828800, RED_ACCENT)
    _add_horizontal_bar(s4, "2. 雇用と経済への影響", 22.3, 35, 2514600, ORANGE_ACCENT)
    _add_horizontal_bar(s4, "3. 人間の自律性の喪失", 21.9, 35, 3200400, PURPLE_ACCENT)

    _add_divider(s4, 4114800, GOLD)

    _add_sub_text(s4, [
        "恩恵は「実体験」として語られた。",
        "懸念の多くは「まだ起きていないが",
        "起きるかもしれない」という予測。",
    ], 4343400, 914400, 12, WARM_GRAY)

    # ── Slide 5: 締め ──
    s5 = _new_slide(prs)
    _add_text_block(s5, [
        "光と影は、",
        "同じ人の中に",
        "ある。",
    ], 1600200, 1600200, 28, WHITE, True)
    _add_divider(s5, 3429000, GOLD)
    _add_text_block(s5, [
        "だからAIは",
        "設計して使う。",
    ], 3657600, 1143000, 28, GOLD, True)
    _add_sub_text(s5, [
        "出典: Anthropic \"81,000 Interviews\" 2026",
    ], 5486400, 457200, 9, WARM_GRAY)

    path = os.path.join(OUTPUT_DIR, "SHORT_NEWS_03_anthropic_light_shade.pptx")
    prs.save(path)
    print(f"Created: {path}")


if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    create_anthropic_freedom()
    create_anthropic_autonomy()
    create_anthropic_light_shade()
    print("\nAll 3 PPTX files created successfully!")
