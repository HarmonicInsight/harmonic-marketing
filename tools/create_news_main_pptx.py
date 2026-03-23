"""
ニュース速報分析シリーズ - 本編用PPTX生成スクリプト
Anthropic 81,000人調査の分析動画（16:9 横型）
HARMONIC insight 標準 Ivory & Gold テーマ (IROVGOLD)
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 定数（16:9） ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================================================
# IROVGOLD カラーパレット (brand/colors.json 準拠)
# ============================================================
BG_PRIMARY = RGBColor(0xFA, 0xF8, 0xF5)      # #FAF8F5 背景
BG_SECONDARY = RGBColor(0xF3, 0xF0, 0xEB)    # #F3F0EB セカンダリ背景
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF カード
BG_HOVER = RGBColor(0xEE, 0xEB, 0xE5)        # #EEEBE5

GOLD_PRIMARY = RGBColor(0xB8, 0x94, 0x2F)    # #B8942F ブランドプライマリ
GOLD_DARK = RGBColor(0x6B, 0x55, 0x18)       # #6B5518
GOLD_LIGHT = RGBColor(0xF0, 0xE6, 0xC8)      # #F0E6C8
GOLD_ACCENT = RGBColor(0xD4, 0xBC, 0x6A)     # #D4BC6A アクセント

TEXT_PRIMARY = RGBColor(0x1C, 0x19, 0x17)     # #1C1917 見出し・本文
TEXT_SECONDARY = RGBColor(0x57, 0x53, 0x4E)   # #57534E 補足テキスト
TEXT_TERTIARY = RGBColor(0xA8, 0xA2, 0x9E)    # #A8A29E キャプション
TEXT_MUTED = RGBColor(0xD6, 0xD3, 0xD1)       # #D6D3D1

BORDER = RGBColor(0xE7, 0xE2, 0xDA)          # #E7E2DA

# セマンティックカラー
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)          # #16A34A
WARNING = RGBColor(0xCA, 0x8A, 0x04)          # #CA8A04
ERROR = RGBColor(0xDC, 0x26, 0x26)            # #DC2626
INFO = RGBColor(0x25, 0x63, 0xEB)             # #2563EB

# アクセントカラー (accent_1〜5)
ACCENT_1 = GOLD_PRIMARY                        # #B8942F
ACCENT_2 = SUCCESS                             # #16A34A
ACCENT_3 = INFO                                # #2563EB
ACCENT_4 = RGBColor(0x7C, 0x3A, 0xED)        # #7C3AED パープル
ACCENT_5 = ERROR                               # #DC2626

WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# フォント (Noto Sans JP 準拠)
FONT_NAME = "Noto Sans JP"
FONT_NAME_FALLBACK = "Yu Gothic"
FONT_NAME_EN = "Segoe UI"

MARGIN_L = Inches(0.8)
CONTENT_W = Inches(11.733)

OUTPUT_DIR = "content/youtube/slides"


# ============================================================
# ヘルパー関数
# ============================================================

def _set_bg(slide, color=BG_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _new_slide(prs, bg_color=BG_PRIMARY):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, bg_color)
    # Top accent bar (Gold)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD_PRIMARY
    bar.line.fill.background()
    return slide


def _add_notes(slide, text):
    """スピーカーノートを追加"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _add_brand_footer(slide):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(7.0), Inches(3), Inches(0.4)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "HARMONIC insight"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.color.rgb = GOLD_PRIMARY


def _add_source_footer(slide, text="出典: Anthropic \"What 81,000 people want from AI\" (2026)"):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6), Inches(7.0), Inches(6.5), Inches(0.4)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_TERTIARY


def _add_title(slide, title, subtitle=None, y=None):
    ty = y or Inches(1.5)
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, ty, CONTENT_W, Inches(1.5)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TEXT_PRIMARY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(12)
        run2 = p2.runs[0]
        run2.font.name = FONT_NAME
        run2.font.size = Pt(18)
        run2.font.color.rgb = TEXT_SECONDARY


def _add_section_header(slide, number, title):
    # Section number circle (Gold)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, MARGIN_L, Inches(2.5), Inches(0.8), Inches(0.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD_PRIMARY
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = str(number)
    r_c.font.name = FONT_NAME_EN
    r_c.font.size = Pt(28)
    r_c.font.bold = True
    r_c.font.color.rgb = WHITE

    # Title
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.4), Inches(9), Inches(1.0)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TEXT_PRIMARY


def _add_horizontal_bar(slide, label, value, max_val, y, color, width_inches=8.0):
    bar_x = Inches(3.5)
    bar_y = Inches(y)
    bar_h = Inches(0.45)
    max_w = Inches(width_inches)
    bar_w = int(max_w * (value / max_val))

    # Label
    lbl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(y - 0.05), Inches(2.5), bar_h
    )
    lbl.fill.background()
    lbl.line.fill.background()
    tf = lbl.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_SECONDARY

    # Bar background
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, max_w, bar_h
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = BG_SECONDARY
    bg_bar.line.fill.background()
    bg_bar.adjustments[0] = 0.15

    # Bar fill
    if bar_w > 0:
        fill_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, Emu(bar_w), bar_h
        )
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = color
        fill_bar.line.fill.background()
        fill_bar.adjustments[0] = 0.15

    # Value
    val_x = Inches(3.5 + width_inches + 0.2)
    val = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, val_x, bar_y, Inches(1), bar_h
    )
    val.fill.background()
    val.line.fill.background()
    tf2 = val.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = f"{value}%"
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = FONT_NAME_EN
    run2.font.size = Pt(16)
    run2.font.bold = True
    run2.font.color.rgb = color


def _add_big_number_center(slide, number, unit, y, color=GOLD_PRIMARY):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(y), Inches(7), Inches(1.5)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run_num = p.add_run()
    run_num.text = str(number)
    run_num.font.name = FONT_NAME_EN
    run_num.font.size = Pt(72)
    run_num.font.bold = True
    run_num.font.color.rgb = color

    run_unit = p.add_run()
    run_unit.text = unit
    run_unit.font.name = FONT_NAME
    run_unit.font.size = Pt(32)
    run_unit.font.bold = True
    run_unit.font.color.rgb = color


def _add_text_block(slide, text, x, y, w, h, font_size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _add_bullet_list(slide, items, x, y, w, h, font_size=16, color=TEXT_PRIMARY):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item if item else " "
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if p.runs:
            run = p.runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size if item else 8)
            run.font.color.rgb = color


def _add_vs_box(slide, val, label, x, y, w, h, val_color):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BG_CARD
    box.line.color.rgb = BORDER
    box.line.width = Pt(1)
    box.adjustments[0] = 0.06

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_val = tf.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    run_v = p_val.add_run()
    run_v.text = val
    run_v.font.name = FONT_NAME_EN
    run_v.font.size = Pt(48)
    run_v.font.bold = True
    run_v.font.color.rgb = val_color

    p_lbl = tf.add_paragraph()
    p_lbl.alignment = PP_ALIGN.CENTER
    run_l = p_lbl.add_run()
    run_l.text = label
    run_l.font.name = FONT_NAME
    run_l.font.size = Pt(14)
    run_l.font.color.rgb = TEXT_SECONDARY


def _add_divider(slide, y, color=GOLD_PRIMARY):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(y), CONTENT_W, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_quote_block(slide, quote, attribution, y):
    # Gold bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y), Inches(0.06), Inches(2.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD_PRIMARY
    bar.line.fill.background()

    _add_text_block(slide, quote, 2.0, y + 0.1, 9.0, 1.5, 20, TEXT_PRIMARY, False)
    _add_text_block(slide, attribution, 2.0, y + 1.5, 9.0, 0.4, 13, TEXT_TERTIARY, False)


def _add_numbered_proposal(slide, num, title, detail, y, color):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = str(num)
    r_c.font.size = Pt(20)
    r_c.font.bold = True
    r_c.font.color.rgb = WHITE
    _add_text_block(slide, title, 1.6, y, 10, 0.35, 18, TEXT_PRIMARY, True)
    _add_text_block(slide, detail, 1.6, y + 0.4, 10, 0.35, 13, TEXT_SECONDARY)


# ============================================================
# スライド生成
# ============================================================
def create_anthropic_81k_main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ━━ Slide 1: タイトル ━━
    s = _new_slide(prs)
    _add_brand_footer(s)
    lbl = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, Inches(0.5), Inches(2.5), Inches(0.45)
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = ERROR
    lbl.line.fill.background()
    lbl.adjustments[0] = 0.25
    tf_l = lbl.text_frame
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_l = tf_l.paragraphs[0]
    p_l.text = "2026年3月18日 発表"
    p_l.alignment = PP_ALIGN.CENTER
    run_l = p_l.runs[0]
    run_l.font.name = FONT_NAME
    run_l.font.size = Pt(13)
    run_l.font.bold = True
    run_l.font.color.rgb = WHITE

    _add_text_block(s, "Anthropic", 3.5, 0.45, 3, 0.5, 14, TEXT_TERTIARY)
    _add_big_number_center(s, "81,000", "人", 1.5, GOLD_PRIMARY)
    _add_text_block(s, "159か国 · 70言語", 3, 3.0, 7, 0.5, 18, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    _add_title(s, "AIに何を求めるか？", "史上最大の質的AI調査を日本のDX視点で分析", Inches(3.8))
    _add_source_footer(s)
    _add_notes(s, """Anthropicが、とんでもない調査を公開しました。

159か国、70言語、81,000人。
「AIに何を求めますか？」という質問を、AIが対話形式で、一人ひとりに深く聞いていった。

チェックボックスのアンケートじゃないんです。
会話で深掘りする質的調査を、8万人規模でやった。史上最大です。

今日はこの調査結果を、日本のDXやビジネスの視点で分析してみます。""")

    # ━━ Slide 2: この調査の何がすごいのか ━━
    s2 = _new_slide(prs)
    _add_brand_footer(s2)
    _add_title(s2, "この調査の何がすごいのか", y=Inches(0.5))
    _add_bullet_list(s2, [
        "従来のAI調査：チェックボックス → 表面的な回答しか得られない",
        "Anthropicの調査：AIが対話で深掘り → 本音が出てくる",
        "",
        "質問：「魔法の杖があったら、AIに何をさせたい？」",
        "→ さらに「なぜそれが欲しいの？」と深層まで掘り下げる",
        "",
        "規模：81,000人 × 159か国 × 70言語",
        "→ 質的調査としては史上最大のスケール",
    ], 0.8, 2.5, 11.5, 4.5, 16, TEXT_PRIMARY)
    _add_source_footer(s2)
    _add_notes(s2, """まず、この調査の何がすごいのかを説明します。

普通のAI調査って、チェックボックスのアンケートなんです。「AIを使っていますか？ はい/いいえ」みたいな。
これだと表面的な答えしか出てこない。

でもAnthropicは違いました。AIが一人ひとりに対話形式で聞いていった。
「魔法の杖があったら、AIに何をさせたい？」
そして「なぜそれが欲しいの？」とさらに深掘りした。

しかもこれを81,000人、159か国、70言語でやった。質的調査としては史上最大です。
量的調査で8万人は珍しくないですが、質的調査でこの規模は前代未聞です。""")

    # ━━ Slide 3: セクション1ヘッダー ━━
    s3 = _new_slide(prs)
    _add_brand_footer(s3)
    _add_section_header(s3, 1, "人々がAIに求めるもの")
    _add_notes(s3, """ではまず、81,000人がAIに何を求めたのかを見ていきます。""")

    # ━━ Slide 4: トップ4バーチャート ━━
    s4 = _new_slide(prs)
    _add_brand_footer(s4)
    _add_title(s4, "AIに求めるもの トップ4", y=Inches(0.4))
    _add_horizontal_bar(s4, "職業的卓越性", 18.8, 25, 2.0, ACCENT_1)
    _add_horizontal_bar(s4, "個人的変容", 13.7, 25, 2.8, ACCENT_2)
    _add_horizontal_bar(s4, "生活管理", 13.5, 25, 3.6, ACCENT_3)
    _add_horizontal_bar(s4, "時間的自由", 11.1, 25, 4.4, ACCENT_4)
    _add_text_block(s4, "上位3つで46%。どれも単純な「生産性向上」には収まらない", 0.8, 5.5, 11, 0.8, 14, TEXT_TERTIARY)
    _add_source_footer(s4)
    _add_notes(s4, """一番多かったのは「職業的卓越性」で19%。
ルーチン業務をAIに任せて、もっと高度な仕事に集中したい。

ここまでは想像通りですよね。
でも、この調査の面白いところは、ここからです。

2位が「個人的変容」14%。成長、メンタルヘルス、自己理解。
3位が「生活管理」14%。認知的負担を減らしたい。
4位が「時間的自由」11%。家族と過ごしたい。

上位3つを合わせると46%。でもどれも、単純な「生産性向上」には収まらない。
これは何を意味しているのか。次のスライドで深掘りします。""")

    # ━━ Slide 5: 深く聞くと変わる ━━
    s5 = _new_slide(prs)
    _add_brand_footer(s5)
    _add_title(s5, "「なぜそれが欲しいの？」と深く聞くと…", y=Inches(0.4))
    _add_bullet_list(s5, [
        "「効率を上げたい」→ 本当は家族と過ごす時間が欲しい",
        "「もっと稼ぎたい」→ 本当は自分の力で独立したい",
        "「業務を自動化したい」→ 本当は頭の中の混乱を整理したい",
    ], 1.0, 2.0, 11.0, 2.5, 18, TEXT_PRIMARY)
    _add_divider(s5, 4.8)
    _add_text_block(s5, "表面的には「生産性」。でも根っこにあるのは「自由」への渇望。", 0.8, 5.2, 11, 0.8, 18, GOLD_PRIMARY, True)
    _add_source_footer(s5)
    _add_notes(s5, """Anthropicは「なぜそれが欲しいの？」と深く聞いていった。
すると、答えが変わるんです。

「効率を上げたい」と言った人。深掘りすると「家族と過ごす時間が欲しい」と言う。
「もっと稼ぎたい」と言った人。深掘りすると「自分の力で独立したい」と言う。
「業務を自動化したい」と言った人。深掘りすると「頭の中の混乱を整理したい」と言う。

人々は「生産性」と答えたけど、本当に求めていたのは「自由」だった。

ここでいう「自由」とは、3つの意味があります。
1つ目は「時間の自由」。家族と過ごしたい、定時に帰りたい。
2つ目は「認知の自由」。頭の中の情報過多から解放されたい。
3つ目は「選択の自由」。自分の意思で人生を決めたい。

この3つが根っこにある。これはすごく大事なポイントです。""")

    # ━━ Slide 6: 具体的な声 ━━
    s6 = _new_slide(prs)
    _add_brand_footer(s6)
    _add_title(s6, "81,000人の「本音」", y=Inches(0.3))
    _add_quote_block(s6,
        "AIのおかげで定時に帰れるようになった。\n子どもの迎えに行ける。",
        "── メキシコ、ソフトウェアエンジニア", 1.5)
    _add_quote_block(s6,
        "AIが認知的負担を引き受けてくれたら、\nそれは家族への「邪魔されない注意」を返してくれる。",
        "── デンマーク、マネージャー", 4.0)
    _add_source_footer(s6)
    _add_notes(s6, """実際の声を紹介します。

メキシコのエンジニアはこう言いました。
「AIのおかげで定時に帰れるようになった。子どもの迎えに行ける。」

デンマークのマネージャーはこう言った。
「AIが認知的負担を本当に引き受けてくれたら、それは何にも代えがたいものを返してくれる。家族への、邪魔されない注意だ。」

この「邪魔されない注意」という表現がすごくリアルじゃないですか。
仕事のことが頭にあると、家族といても心ここにあらずになる。
AIが仕事の認知負荷を下げてくれたら、目の前の人に集中できるようになる。

生産性向上って、結局そういうことなんです。""")

    # ━━ Slide 7: KPIを変えろ ━━
    s7 = _new_slide(prs)
    _add_brand_footer(s7)
    _add_title(s7, "AI導入の目標設定を変える", y=Inches(0.4))
    _add_vs_box(s7, "×", "処理速度が何%上がったか", 1.0, 2.0, 5.0, 2.5, ERROR)
    _add_vs_box(s7, "○", "社員が本来の仕事に\n集中できているか", 7.0, 2.0, 5.0, 2.5, SUCCESS)
    _add_text_block(s7, "生産性は手段。目的は「人が人生を取り戻すこと」。", 0.8, 5.5, 11, 0.8, 18, GOLD_PRIMARY, True)
    _add_source_footer(s7)
    _add_notes(s7, """これ、すごく大事なポイントです。
AIを導入する時の目標設定が変わります。

今までのKPI：「処理速度が何%上がったか」「コストが何%下がったか」
これは間違いではないけど、本質を見失っている。

本当に測るべきは「社員が本来やるべき仕事に集中できているか」。

処理速度が2倍になっても、社員が疲弊していたら意味がない。
逆に、処理速度は変わらなくても、社員が創造的な仕事に時間を使えるようになったら、それは大きな成功です。

生産性は手段であって、目的じゃない。目的は「人が人生を取り戻すこと」。""")

    # ━━ Slide 8: セクション2ヘッダー ━━
    s8 = _new_slide(prs)
    _add_brand_footer(s8)
    _add_section_header(s8, 2, "衝撃のデータ：自律性ギャップ")
    _add_notes(s8, """次に、僕がこの調査で一番注目したデータを紹介します。""")

    # ━━ Slide 9: 47% vs 14% ━━
    s9 = _new_slide(prs)
    _add_brand_footer(s9)
    _add_title(s9, "AIで経済的恩恵を感じている人の割合", y=Inches(0.4))
    _add_vs_box(s9, "47%", "独立事業者\nフリーランス", 1.0, 2.0, 5.0, 2.8, SUCCESS)
    _add_vs_box(s9, "14%", "企業の従業員", 7.0, 2.0, 5.0, 2.8, ERROR)
    vs = s9.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(6.0), Inches(2.8), Inches(1.0), Inches(1.0)
    )
    vs.fill.solid()
    vs.fill.fore_color.rgb = ERROR
    vs.line.fill.background()
    tf_vs = vs.text_frame
    tf_vs.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_vs = tf_vs.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    run_vs = p_vs.add_run()
    run_vs.text = "VS"
    run_vs.font.name = FONT_NAME_EN
    run_vs.font.size = Pt(18)
    run_vs.font.bold = True
    run_vs.font.color.rgb = WHITE
    _add_text_block(s9, "3.4倍の差", 4.5, 5.3, 4, 0.6, 28, GOLD_PRIMARY, True, PP_ALIGN.CENTER)
    _add_source_footer(s9)
    _add_notes(s9, """AIで経済的な恩恵を感じている人の割合。

独立事業者、フリーランス：47%。
企業の従業員：14%。

3倍以上の差です。

同じAIを使っているのに、なぜここまで差がつくのか。
フリーランスは「こう使ったら便利だ」と思ったら、その場で試せる。
自分で判断して、すぐに実行できる。

でも会社員は違う。「まずセキュリティ部門に確認します」「来期の予算で検討」。
この承認プロセスの間に、フリーランスはもう成果を出している。""")

    # ━━ Slide 10: 副業持ち ━━
    s10 = _new_slide(prs)
    _add_brand_footer(s10)
    _add_title(s10, "さらに興味深いデータ", y=Inches(0.4))
    _add_text_block(s10, "副業を持っている会社員", 3.0, 1.8, 7, 0.6, 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    _add_big_number_center(s10, "58", "%", 2.5, WARNING)
    _add_divider(s10, 4.3)
    _add_bullet_list(s10, [
        "会社員でも「自分の裁量で使える場」があれば恩恵を受けられる",
        "つまりAIの恩恵は、AIの性能ではなく「使い方の自由度」で決まる",
    ], 1.0, 4.6, 11.0, 2.0, 16, TEXT_PRIMARY)
    _add_source_footer(s10)
    _add_notes(s10, """さらに面白いデータがあります。

副業を持っている会社員は、58%がAIの経済的恩恵を感じている。

会社では14%なのに、副業では58%。同じ人が、です。
何が違うか。副業では自分の裁量でAIを使えるからです。

つまり、AIの恩恵を受けられるかどうかは、AIの性能の問題じゃない。
自分の裁量で使えるかどうか、なんです。

これは組織にとって非常に重要な示唆です。
「良いAIツールを入れれば生産性が上がる」は幻想。
社員に使い方の自由を渡さなければ、どんなに優れたAIも宝の持ち腐れです。""")

    # ━━ Slide 11: なぜ差が生まれるのか ━━
    s11 = _new_slide(prs)
    _add_brand_footer(s11)
    _add_title(s11, "なぜ3.4倍の差が生まれるのか", y=Inches(0.4))
    _add_bullet_list(s11, [
        "フリーランス：自分で判断 → すぐ実行 → すぐ成果",
        "",
        "会社員：",
        "  「まずセキュリティ部門に確認します」",
        "  「AIツール導入は来期の予算で検討」",
        "  「個人情報の取り扱いが…承認プロセスが…」",
        "",
        "→ その間に、フリーランスはAIで3日で終わらせている",
    ], 0.8, 1.8, 11.5, 4.0, 16, TEXT_PRIMARY)
    _add_text_block(s11, "問題はAIの性能ではない。自律性を渡しているかどうか。", 0.8, 5.8, 11, 0.8, 18, GOLD_PRIMARY, True)
    _add_source_footer(s11)
    _add_notes(s11, """具体的に何が起きているかというと、

フリーランスの場合：
「この提案書、AIに下書きさせよう」→ 5分で実行 → 1時間で完成。

会社員の場合：
「AI使いたいんですけど」→「まずセキュリティ部門に確認して」
→ 2週間後「承認されました、でもこのツールは使えません」
→「じゃあ別のツールを申請します」→ また2週間…

この間にフリーランスは10件の仕事を終わらせている。

問題はAIの性能じゃない。自律性を渡しているかどうか。
これは個人の能力の問題でもない。組織の仕組みの問題です。""")

    # ━━ Slide 12: 日本企業のDX ━━
    s12 = _new_slide(prs)
    _add_brand_footer(s12)
    _add_title(s12, "日本企業のDXに置き換えると", y=Inches(0.4))
    _add_text_block(s12, "日本企業の", 3.5, 1.8, 6, 0.5, 20, TEXT_PRIMARY, False, PP_ALIGN.CENTER)
    _add_big_number_center(s12, "79.3", "%", 2.2, ERROR)
    _add_text_block(s12, "がDXの最大課題は「人材不足」と回答", 2.5, 3.6, 8, 0.5, 18, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    _add_divider(s12, 4.3)
    _add_text_block(s12, "本当の問題は、人材がいないことではなく\nいる人材に裁量を渡していないこと", 1.5, 4.6, 10, 1.5, 22, GOLD_PRIMARY, True, PP_ALIGN.CENTER)
    _add_source_footer(s12, "出典: 企業IT動向調査2025（JUAS）/ Anthropic 81,000 Interviews")
    _add_notes(s12, """ここに日本企業のDXの課題が凝縮されています。

別の調査では、日本企業の79.3%がDXの最大課題は「人材不足」だと答えています。

でも、本当にそうでしょうか。

47% vs 14%のデータを思い出してください。
同じ人が、会社では恩恵を受けられず、副業では恩恵を受けている。
人材がいないんじゃなくて、いる人材に裁量を渡していない。

「まず上に確認します」「前例がないので」「規程では認められていない」
この3つの言葉が、日本企業のAI活用を止めている。

自律性を渡さない組織は、AIの恩恵を受けられない。
これはかなり強いメッセージだと思います。""")

    # ━━ Slide 13: セクション3ヘッダー ━━
    s13 = _new_slide(prs)
    _add_brand_footer(s13)
    _add_section_header(s13, 3, "光と影 ── 同じ人の中に共存する")
    _add_notes(s13, """この調査で最も深い発見が、この「光と影」です。

AIの利点と懸念は、賛成派と反対派に分かれるのではない。
同じ人の中に共存している。これがこの調査のもっとも重要なメッセージです。""")

    # ━━ Slide 14: 光と影データ ━━
    s14 = _new_slide(prs)
    _add_brand_footer(s14)
    _add_title(s14, "AIの利点と懸念は、別の人が感じているのではない", y=Inches(0.3))
    _add_text_block(s14, "同じ人の中に共存している", 0.8, 1.2, 11, 0.6, 22, GOLD_PRIMARY, True)
    _add_vs_box(s14, "3倍", "AIに感情的に助けられた人\n→ 依存を恐れる確率", 0.8, 2.3, 5.5, 2.0, ERROR)
    _add_vs_box(s14, "2.5倍", "教師が学生の\n思考力低下を目撃する率", 6.8, 2.3, 5.5, 2.0, WARNING)
    _add_text_block(s14, "恩恵は「実体験」として語られた（81%が「もう恩恵を受けている」）", 0.8, 4.8, 11, 0.5, 14, TEXT_SECONDARY)
    _add_text_block(s14, "懸念の多くは「まだ起きていないが、起きるかもしれない」という予測", 0.8, 5.3, 11, 0.5, 14, TEXT_TERTIARY)
    _add_text_block(s14, "ただし例外：フリーランス創作者（23%恩恵 vs 17%脅威 → ほぼ拮抗、現実化）", 0.8, 5.8, 11, 0.5, 14, WARNING)
    _add_source_footer(s14)
    _add_notes(s14, """データを見てください。

AIに感情的に助けられたと感じている人。
その人たちが、AIへの依存を恐れる確率は3倍高かった。

AIで学習が加速したと感じている人。
その人たちが、認知萎縮を心配する確率も高い。

教育の現場では、教師が学生の思考力低下を平均の2.5倍の頻度で目撃している。

ここが重要なんですが、恩恵は実体験として語られていた。81%が「もう恩恵を受けている」と。
一方、懸念の多くは「まだ起きていないけど、起きるかもしれない」という予測だった。

ただし例外があります。フリーランスの創作者。
23%が恩恵を感じ、17%が脅威を感じている。ほぼ拮抗。
これはもう予測じゃなくて、現実です。""")

    # ━━ Slide 15: 具体的な声（光と影） ━━
    s15 = _new_slide(prs)
    _add_brand_footer(s15)
    _add_title(s15, "使っている人ほど、恐れている", y=Inches(0.3))
    _add_quote_block(s15,
        "AIで契約書レビューの時間を節約している。\nでも同時に恐れている。自分で読む力を失っているのではないか。",
        "── イスラエル、弁護士", 1.5)
    _add_quote_block(s15,
        "翻訳AIのおかげで海外の論文が読めるようになった。\nでも英語で考える力が確実に落ちている。半年前の自分の方が読めた。",
        "── 日本のビジネスパーソンに置き換えると…", 4.0)
    _add_source_footer(s15)
    _add_notes(s15, """具体的な声を紹介します。

イスラエルの弁護士はこう言いました。
「AIで契約書レビューの時間を節約している。でも同時に恐れている。自分で読む力を失っているのではないか。」

これ、日本のビジネスパーソンにも当てはまりませんか？

例えば、翻訳AI。おかげで海外の論文や記事がサクサク読めるようになった。
でも、英語で考える力が確実に落ちている。半年前の自分の方が英文を読めた気がする。

あるいは、ChatGPTで企画書のたたき台を作るようになった。
でも、自分でゼロから考える力が弱くなっている気がする。

これが「光と影は同じ人の中にある」ということの意味です。
便利だと感じている人こそが、自分の能力低下を一番敏感に感じている。""")

    # ━━ Slide 16: 懸念トップ3 ━━
    s16 = _new_slide(prs)
    _add_brand_footer(s16)
    _add_title(s16, "81,000人が挙げた懸念 トップ3", y=Inches(0.4))
    _add_horizontal_bar(s16, "信頼性（ハルシネーション）", 26.7, 35, 2.2, ERROR)
    _add_horizontal_bar(s16, "雇用と経済への影響", 22.3, 35, 3.2, WARNING)
    _add_horizontal_bar(s16, "人間の自律性の喪失", 21.9, 35, 4.2, ACCENT_4)
    _add_divider(s16, 5.3)
    _add_text_block(s16, "11%は「懸念なし」と回答。中・低所得国に多い（アフリカ18%）", 0.8, 5.6, 11, 0.5, 14, TEXT_SECONDARY)
    _add_text_block(s16, "先進国：「今あるものを奪う脅威」 ↔ 途上国：「今までなかったものを与える希望」", 0.8, 6.1, 11, 0.5, 14, TEXT_TERTIARY)
    _add_source_footer(s16)
    _add_notes(s16, """懸念のトップ3も紹介します。

1位：信頼性の問題。ハルシネーション。27%。
AIが嘘をつく。これが一番の懸念です。

2位：雇用と経済への影響。22%。
自分の仕事がなくなるのではないか。

3位：人間の自律性の喪失。22%。
AIに依存しすぎて、自分で判断する力を失うのではないか。

面白いのは、11%の人は「懸念なし」と答えていること。
しかもこれは中・低所得国に多い。アフリカでは18%が懸念なし。

なぜか。
先進国では、AIは「今あるものを奪う脅威」。
途上国では、AIは「今までなかったものを与えてくれる希望」。
同じテクノロジーが、立場で全く違って見えるんです。""")

    # ━━ Slide 17: セクション4ヘッダー ━━
    s17 = _new_slide(prs)
    _add_brand_footer(s17)
    _add_section_header(s17, 4, "日本への示唆 ── 東アジアのデータから")
    _add_notes(s17, """最後に、この調査を日本のDXにどう活かすか。
東アジアのデータも出ていますので、そこから見ていきます。""")

    # ━━ Slide 18: 東アジアデータ ━━
    s18 = _new_slide(prs)
    _add_brand_footer(s18)
    _add_title(s18, "東アジア（日本含む）の特徴", y=Inches(0.4))
    _add_horizontal_bar(s18, "個人的変容への期待", 19, 25, 2.2, ACCENT_3)
    _add_horizontal_bar(s18, "経済的独立への期待", 15, 25, 3.0, ACCENT_2)
    _add_horizontal_bar(s18, "認知萎縮への懸念", 18, 25, 3.8, ERROR)
    _add_horizontal_bar(s18, "創造性喪失への懸念", 13, 25, 4.6, WARNING)
    _add_text_block(s18, "AIに自己変革を最も期待し、自分が弱くなることを最も恐れている地域", 0.8, 5.8, 11, 0.8, 16, GOLD_PRIMARY, True)
    _add_source_footer(s18)
    _add_notes(s18, """東アジアのデータを見てみましょう。

個人的変容への期待：19%で全地域最高。
経済的独立への期待：15%で全地域最高。
認知萎縮への懸念：18%。
意味と創造性の喪失への懸念：13%。

つまり東アジア、特に日本は、
AIに自己変革を期待しつつ、自分が弱くなることを最も恐れている地域なんです。

「AIで成長したい。でもAIに頼りすぎて弱くなりたくない。」
この二律背反が、日本のAI活用を難しくしている根本原因かもしれません。

だからこそ、「とりあえずAIを使え」ではダメなんです。
意識的に設計して使う必要がある。""")

    # ━━ Slide 19: 3つの提言 ━━
    s19 = _new_slide(prs)
    _add_brand_footer(s19)
    _add_title(s19, "81,000人の声から、3つの提言", y=Inches(0.3))
    _add_numbered_proposal(s19, 1,
        "目標を「効率化」から「自由の創出」に変える",
        "処理速度ではなく、人が本来やるべき仕事に集中できているかを測る",
        1.8, ACCENT_3)
    _add_numbered_proposal(s19, 2,
        "自律性を渡す ── 承認を減らし、裁量を増やす",
        "47% vs 14%のデータが全て。使い方を現場に委ねる",
        3.2, ACCENT_2)
    _add_numbered_proposal(s19, 3,
        "AIの使い方を「設計する」── とりあえず使えはダメ",
        "何をAIに任せ、何を自分の力として残すか。意識的に設計する",
        4.6, ACCENT_1)
    _add_source_footer(s19)
    _add_notes(s19, """ここから3つ、提言します。

1つ目。AIの導入目標を「効率化」から「自由の創出」に変える。
処理速度ではなく、人が本来やるべき仕事に集中できているかを測る。
KPIの設計からやり直す必要があります。

2つ目。自律性を渡す。
47% vs 14%のデータが全てです。
AIを活用したいなら、使い方を現場に委ねる。承認を減らし、裁量を増やす。
「まず使わせる。問題が起きたら対処する」の順番に変える。

3つ目。AIの使い方を「設計する」。
光と影は同じ人の中にある。だから「とりあえず使え」ではダメ。
何をAIに任せ、何を自分の力として残すか。意識的に設計する必要がある。

例えば「翻訳はAIに任せるが、週に1回は自分で英文を読む時間を取る」。
こういうルールを自分で決める。""")

    # ━━ Slide 20: 締め ━━
    s20 = _new_slide(prs)
    _add_brand_footer(s20)
    _add_text_block(s20, "AIは道具であり、同時に鏡である。", 1.5, 2.0, 10, 1.0, 28, TEXT_PRIMARY, True, PP_ALIGN.CENTER)
    _add_divider(s20, 3.3)
    _add_text_block(s20, "自分が何を求め、何を恐れているか。\nそれが映し出される。", 1.5, 3.6, 10, 1.2, 22, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    _add_text_block(s20, "だからこそ、AIの前に\n自分が本当に欲しいものを明確にすること。", 1.5, 5.0, 10, 1.2, 22, GOLD_PRIMARY, True, PP_ALIGN.CENTER)
    _add_source_footer(s20)
    _add_notes(s20, """81,000人の声が教えてくれたのは、
AIは道具であると同時に、鏡だということです。

自分が何を求め、何を恐れているか。それが映し出される。

生産性を求めた人は、実は自由を求めていた。
便利さを感じた人は、同時に自分の能力低下を恐れていた。
会社で使えない人は、組織の自律性の問題を映し出していた。

だからこそ、AIの前に、自分が本当に欲しいものを明確にすることが大事なんです。""")

    # ━━ Slide 21: エンディング ━━
    s21 = _new_slide(prs)
    _add_brand_footer(s21)
    _add_text_block(s21, "HARMONIC insight", 2.5, 2.0, 8, 1.0, 36, GOLD_PRIMARY, True, PP_ALIGN.CENTER)
    _add_text_block(s21, "AI時代のビジネス戦略を発信しています", 2.5, 3.2, 8, 0.6, 18, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    _add_divider(s21, 4.2)
    _add_bullet_list(s21, [
        "調査全文: anthropic.com/81k-interviews",
        "YouTube: チャンネル登録お願いします",
        "note: note.com/harmonic_insight",
    ], 3.0, 4.6, 7, 2.0, 14, TEXT_SECONDARY)
    _add_notes(s21, """今回の内容が参考になったら、チャンネル登録をお願いします。
Anthropicの調査全文はリンクを概要欄に貼っておきます。

DXやAI活用の実践的な話を発信しています。
次回もお楽しみに。""")

    # Save
    path = os.path.join(OUTPUT_DIR, "VID-NEWS_01_anthropic_81k_analysis.pptx")
    prs.save(path)
    print(f"Created: {path}")


if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    create_anthropic_81k_main()
    print("\nMain PPTX created successfully!")
