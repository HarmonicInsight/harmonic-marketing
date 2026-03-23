"""
AI時代の差別化 - Shorts用PPTX生成スクリプト
テーマ：「機能の時代は終わる。これからは相性で選ばれる」
既存デザイン（ダーク背景 + ゴールドアクセント）を踏襲。
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 定数 ──
SLIDE_W = Emu(4114800)  # 4.5 inches
SLIDE_H = Emu(7315200)  # 8.0 inches

# カラーパレット
BG_DARK = RGBColor(0x1C, 0x19, 0x17)
GOLD = RGBColor(0xB8, 0x94, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM_GRAY = RGBColor(0x57, 0x53, 0x4E)
LIGHT_GRAY = RGBColor(0xA8, 0xA2, 0x9E)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)
BLUE_ACCENT = RGBColor(0x42, 0xA5, 0xF5)
GREEN_ACCENT = RGBColor(0x66, 0xBB, 0x6A)
ORANGE_ACCENT = RGBColor(0xFF, 0xA7, 0x26)
TEAL_ACCENT = RGBColor(0x26, 0xC6, 0xDA)
DARK_SURFACE = RGBColor(0x2C, 0x27, 0x23)
PURPLE_ACCENT = RGBColor(0xAB, 0x47, 0xBC)

FONT_NAME = "Hiragino Sans"
FONT_NAME_FALLBACK = "Yu Gothic"

MARGIN_X = Emu(274320)
CONTENT_W = Emu(3566160)
TOP_BAR_H = Emu(54864)
BRAND_Y = Emu(6675120)
BRAND_H = Emu(365760)

OUTPUT_DIR = "content/youtube/shorts"


def _set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_top_bar(slide, color=GOLD):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TOP_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_brand(slide):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, BRAND_Y, CONTENT_W, BRAND_H
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HARMONIC insight"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(9)
    run.font.color.rgb = GOLD


def _new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide)
    _add_top_bar(slide)
    _add_brand(slide)
    return slide


def _add_text_block(slide, texts, y, h, font_size, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y), CONTENT_W, Emu(h)
    )
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_sub_text(slide, texts, y, h, font_size=12, color=WARM_GRAY):
    _add_text_block(slide, texts, y, h, font_size, color, bold=False)


def _add_divider(slide, y, color=GOLD):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Emu(y), CONTENT_W, Emu(18288)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_icon_text_row(slide, icon_char, text, y, icon_color=GOLD, text_color=WHITE):
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, MARGIN_X, Emu(y), Emu(320040), Emu(320040)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = DARK_SURFACE
    icon_box.line.fill.background()
    tf_i = icon_box.text_frame
    tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_i = tf_i.paragraphs[0]
    p_i.alignment = PP_ALIGN.CENTER
    run_i = p_i.add_run()
    run_i.text = icon_char
    run_i.font.name = FONT_NAME
    run_i.font.size = Pt(14)
    run_i.font.bold = True
    run_i.font.color.rgb = icon_color

    txt_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(MARGIN_X + 388620), Emu(y), Emu(CONTENT_W - 388620), Emu(320040)
    )
    txt_box.fill.background()
    txt_box.line.fill.background()
    tf_t = txt_box.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]
    p_t.text = text
    run_t = p_t.runs[0]
    run_t.font.name = FONT_NAME
    run_t.font.size = Pt(13)
    run_t.font.color.rgb = text_color


def _add_arrow_flow(slide, items, start_y, spacing, color=GOLD):
    """矢印付きフローチャート（縦並び）"""
    for i, (text, item_color) in enumerate(items):
        y = start_y + i * spacing

        # ボックス
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, Emu(y),
            CONTENT_W, Emu(365760)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_SURFACE
        box.line.color.rgb = item_color
        box.line.width = Pt(1.5)
        box.adjustments[0] = 0.15

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = item_color

        # 矢印（最後のアイテム以外）
        if i < len(items) - 1:
            arrow_y = y + 365760 + int((spacing - 365760) * 0.2)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Emu(int(MARGIN_X + CONTENT_W / 2 - 91440)),
                Emu(arrow_y), Emu(182880), Emu(int((spacing - 365760) * 0.6))
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = WARM_GRAY
            arrow.line.fill.background()


def _add_comparison_block(slide, before_title, before_items, after_title, after_items, y):
    """従来 vs AI時代 の比較ブロック"""
    half_w = int(CONTENT_W / 2) - Emu(45720)
    block_h = Emu(1600200)

    # 従来（左）
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, Emu(y), Emu(half_w), block_h
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = DARK_SURFACE
    left_box.line.fill.background()
    left_box.adjustments[0] = 0.08

    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.vertical_anchor = MSO_ANCHOR.TOP

    p_title = tf_l.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    p_title.space_after = Pt(8)
    run_t = p_title.add_run()
    run_t.text = before_title
    run_t.font.name = FONT_NAME
    run_t.font.size = Pt(12)
    run_t.font.bold = True
    run_t.font.color.rgb = RED_ACCENT

    for item in before_items:
        p = tf_l.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = item
        r.font.name = FONT_NAME
        r.font.size = Pt(11)
        r.font.color.rgb = LIGHT_GRAY

    # AI時代（右）
    right_x = Emu(int(MARGIN_X + half_w + 91440))
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Emu(y), Emu(half_w), block_h
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = DARK_SURFACE
    right_box.line.color.rgb = GOLD
    right_box.line.width = Pt(1.5)
    right_box.adjustments[0] = 0.08

    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.vertical_anchor = MSO_ANCHOR.TOP

    p_title2 = tf_r.paragraphs[0]
    p_title2.alignment = PP_ALIGN.CENTER
    p_title2.space_after = Pt(8)
    run_t2 = p_title2.add_run()
    run_t2.text = after_title
    run_t2.font.name = FONT_NAME
    run_t2.font.size = Pt(12)
    run_t2.font.bold = True
    run_t2.font.color.rgb = GOLD

    for item in after_items:
        p = tf_r.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = item
        r.font.name = FONT_NAME
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE


# ============================================================
# AI時代の差別化：機能の時代は終わる
# ============================================================
def create_ai_era_design_focus():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: タイトル ──
    s1 = _new_slide(prs)
    _add_text_block(s1, [
        "AI時代、",
        "「機能」では",
        "もう勝てない",
    ], 1143000, 1600200, 28, WHITE, True)
    _add_divider(s1, 2971800, GOLD)
    _add_text_block(s1, [
        "差がつくのは",
        "「相性」だけ",
    ], 3200400, 1143000, 28, GOLD, True)
    _add_sub_text(s1, [
        "機能格差ゼロ時代のシステム選び",
    ], 4571040, 457200, 13, WARM_GRAY)

    # ── Slide 2: フロー図 - 機能差→性能差→体験差 ──
    s2 = _new_slide(prs)
    _add_text_block(s2, [
        "競争の行方",
    ], 457200, 548640, 22, GOLD, True)
    _add_sub_text(s2, [
        "差別化ポイントの変遷",
    ], 914400, 365760, 13, LIGHT_GRAY)

    _add_arrow_flow(s2, [
        ("機能差は消える", RED_ACCENT),
        ("性能差も縮む", ORANGE_ACCENT),
        ("体験・思想・信頼・好み", GOLD),
    ], 1554480, 640080)

    _add_sub_text(s2, [
        "基盤モデルが共通化",
        "ベストプラクティスが広がる",
        "UI/UXパターンが標準化",
        "各社とも同じ最適化をする",
    ], 3886200, 1143000, 11, WARM_GRAY)

    # ── Slide 3: なぜ機能差が消えるか ──
    s3 = _new_slide(prs)
    _add_text_block(s3, [
        "なぜ機能では",
        "差がつかないか",
    ], 822960, 914400, 22, WHITE, True)
    _add_divider(s3, 1828800, GOLD)

    _add_icon_text_row(s3, "1", "すぐに模倣される", 2114550, RED_ACCENT, WHITE)
    _add_icon_text_row(s3, "2", "統合・API化される", 2571750, ORANGE_ACCENT, WHITE)
    _add_icon_text_row(s3, "3", "同等機能が並ぶ", 3028950, BLUE_ACCENT, WHITE)

    _add_divider(s3, 3657600, GOLD)
    _add_text_block(s3, [
        "独自機能という状態は",
        "長く続かない",
    ], 3886200, 914400, 18, GOLD, True)

    # ── Slide 4: 最後に残る4つの差 ──
    s4 = _new_slide(prs)
    _add_text_block(s4, [
        "最後に残る差",
    ], 457200, 548640, 22, GOLD, True)

    _add_icon_text_row(s4, "1", "性格 ─ 慎重か、大胆か", 1143000, BLUE_ACCENT, WHITE)
    _add_icon_text_row(s4, "2", "フィット感 ─ 自然に使えるか", 1600200, GREEN_ACCENT, WHITE)
    _add_icon_text_row(s4, "3", "信頼 ─ 間違えた時どう振る舞うか", 2057400, ORANGE_ACCENT, WHITE)
    _add_icon_text_row(s4, "4", "好き ─ 使っていて気分がよいか", 2514600, TEAL_ACCENT, WHITE)

    _add_divider(s4, 3200400, GOLD)
    _add_text_block(s4, [
        "能力差が極端にないなら",
        "最後は「相性」で決まる",
    ], 3429000, 914400, 18, WHITE, True)
    _add_sub_text(s4, [
        "車や家電を選ぶというより、",
        "「人を選ぶ」感覚に近い",
    ], 4571040, 640080, 12, WARM_GRAY)

    # ── Slide 5: 従来 vs AI時代 比較 ──
    s5 = _new_slide(prs)
    _add_text_block(s5, [
        "何が変わるのか",
    ], 457200, 548640, 22, GOLD, True)

    _add_comparison_block(
        s5,
        "従来",
        ["機能が多い方が勝つ", "できることが多い", "方が勝つ"],
        "AI時代",
        ["機能は揃う", "性能も近づく", "「好きか」「任せ", "られるか」で決まる"],
        1371600,
    )

    _add_divider(s5, 3200400, GOLD)
    _add_text_block(s5, [
        "作る側が設計すべきは",
        "「機能一覧」ではなく",
    ], 3429000, 914400, 16, WHITE, True)
    _add_text_block(s5, [
        "どんな人格と体験と",
        "仕事観を持つシステムか",
    ], 4343400, 731520, 16, GOLD, True)

    # ── Slide 6: 結論 ──
    s6 = _new_slide(prs)
    _add_text_block(s6, [
        "機能の時代は",
        "終わる。",
    ], 1371600, 1143000, 28, WHITE, True)
    _add_divider(s6, 2743200, GOLD)
    _add_text_block(s6, [
        "これからは",
        "「相性」で",
        "選ばれる。",
    ], 2971800, 1600200, 28, GOLD, True)
    _add_sub_text(s6, [
        "AI時代のシステム差別化は、",
        "機能ではなく「体験の哲学」になる。",
    ], 5029200, 640080, 11, WARM_GRAY)

    path = os.path.join(OUTPUT_DIR, "SHORT_ai_era_design_focus.pptx")
    prs.save(path)
    print(f"Created: {path}")


if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    create_ai_era_design_focus()
    print("\nPPTX file created successfully!")
